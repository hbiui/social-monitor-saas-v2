#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
豆包 AI PPT 生成流水线 v4  —  Doubao Professional Style (参考豆包PPTX设计语言)
- 深色背景: 111827 (Tailwind gray-900)
- 白色/深色卡片双模式
- 三色品牌系统: 蓝/绿/琥珀
- 放射渐变装饰圆球 (封面/结尾)
- 无背景图依赖的全矢量设计
"""

import sys, json, os, time, base64, urllib.request, urllib.error, io, re, math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

try:
    from PIL import Image as PILImage
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

ARK_BASE  = 'https://ark.cn-beijing.volces.com/api/v3'
CHAT_URL  = f'{ARK_BASE}/chat/completions'
IMG_URL   = f'{ARK_BASE}/images/generations'
IMG_MODEL = 'doubao-seedream-5-0-260128'

# ── Doubao Reference Color System (exact match from reference PPTX) ────────────
C_BG     = RGBColor(0x11, 0x18, 0x27)  # Tailwind gray-900 — slide background
C_CARD   = RGBColor(0x1F, 0x29, 0x37)  # gray-800 — card background (white variant)
C_CARD2  = RGBColor(0x2A, 0x2D, 0x33)  # — dark card
C_CARD3  = RGBColor(0x2C, 0x31, 0x3A)  # — darker card (recommendations)
C_CARD4  = RGBColor(0x23, 0x27, 0x2E)  # — darkest card
C_BLUE   = RGBColor(0x3B, 0x82, 0xF6)  # blue-500 — primary accent
C_BLUE_L = RGBColor(0x93, 0xC5, 0xFD)  # blue-300 — label text
C_BLUE_M = RGBColor(0x60, 0xA5, 0xFA)  # blue-400 — heading text
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)  # emerald-500
C_GREEN_L= RGBColor(0x6E, 0xE7, 0xB7)  # emerald-300
C_GREEN_M= RGBColor(0x34, 0xD3, 0x99)  # emerald-400
C_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)  # amber-500
C_AMBER_L= RGBColor(0xFC, 0xD3, 0x4D)  # amber-300
C_AMBER_M= RGBColor(0xFB, 0xBF, 0x24)  # amber-400
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_T1     = RGBColor(0xE5, 0xE7, 0xEB)  # gray-200 — primary text
C_T2     = RGBColor(0xD1, 0xD5, 0xDB)  # gray-300 — body text
C_T3     = RGBColor(0x9C, 0xA3, 0xAF)  # gray-400 — muted text
C_T4     = RGBColor(0x6B, 0x72, 0x80)  # gray-500 — dim text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Brand color triples: (header, label_light, text_accent)
BRAND_COLORS = [
    (C_BLUE,  C_BLUE_L,  C_BLUE_M),   # blue
    (C_GREEN, C_GREEN_L, C_GREEN_M),  # green
    (C_AMBER, C_AMBER_L, C_AMBER_M),  # amber
    (C_BLUE,  C_BLUE_L,  C_BLUE_M),   # loop
    (C_GREEN, C_GREEN_L, C_GREEN_M),
]


# ═══════════════════════════════════════════════════════════════════════════════
# Status
# ═══════════════════════════════════════════════════════════════════════════════

def write_status(path, status, step, message, progress, url=None, error=None):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump({'status':status,'step':step,'message':message,
                   'progress':progress,'url':url,'error':error,
                   'updated':time.time()}, f, ensure_ascii=False)


# ═══════════════════════════════════════════════════════════════════════════════
# HTTP helpers
# ═══════════════════════════════════════════════════════════════════════════════

def make_opener(proxy_url=''):
    if proxy_url:
        h = urllib.request.ProxyHandler({'http':proxy_url,'https':proxy_url})
        return urllib.request.build_opener(h)
    return urllib.request.build_opener()

def http_post(url, headers, body_bytes, proxy_url='', timeout=180):
    req = urllib.request.Request(url, data=body_bytes, headers=headers, method='POST')
    opener = make_opener(proxy_url)
    with opener.open(req, timeout=timeout) as r: return r.read()

def doubao_chat(api_key, text_model, messages, proxy_url='', max_tokens=12000):
    body = {'model':text_model,'messages':messages,'max_tokens':max_tokens,'temperature':0.3}
    raw = http_post(CHAT_URL,
                    {'Content-Type':'application/json','Authorization':f'Bearer {api_key}'},
                    json.dumps(body).encode(), proxy_url, 300)
    return json.loads(raw)['choices'][0]['message']['content']

# Seedream supports these 16:9 sizes (try in order)
# Seedream size params: API only accepts quality keywords or valid aspect ratios
# '2K' → ~2048px widescreen  '1K' → ~1024px  '16:9' → explicit ratio (newer API support)
# Do NOT use pixel strings like '1280x720' — these are rejected with HTTP 400
SEEDREAM_SIZES_16_9 = ['1920x1080', '1280x720', '2K', '1K']

def seedream_generate(api_key, prompt, proxy_url=''):
    """
    Call doubao-seedream-5-0-260128 via ARK images/generations endpoint.

    ROOT CAUSE FIXED (v14):
    - Seedream API only accepts quality keywords like "2K", "1K" for size,
      NOT raw pixel strings like "1280x720" (those return HTTP 400).
    - Also required: stream=false, response_format="url", watermark=false
    - Try 16:9 ratio first (widescreen for PPT), then fallback to "2K"
    """
    if not api_key:
        print('[Seedream] ❌ No API key provided', file=sys.stderr)
        return None

    # PPT-optimised size candidates for Seedream 5.0:
    # '1920x1080' may work for explicit 16:9, fallback to quality keywords
    SIZE_CANDIDATES = [
        {'size': '1920x1080'},          # 16:9 explicit — preferred for PPT
        {'size': '2K'},                 # ~2048px, auto aspect — fallback 1
        {'size': '1K'},                 # ~1024px, auto aspect — fallback 2
        {},                             # omit size entirely → model default
    ]

    headers = {
        'Content-Type':  'application/json',
        'Authorization': f'Bearer {api_key}',
    }

    last_err = None
    for size_params in SIZE_CANDIDATES:
        body = {
            'model':  IMG_MODEL,
            'prompt': prompt[:2000],
            'n':      1,
            'stream': False,               # required for Seedream 5.0
            'response_format': 'url',
            'watermark': False,
            'sequential_image_generation': 'disabled',
        }
        body.update(size_params)   # add size key only if non-empty

        size_label = size_params.get('size','default')
        try:
            print(f'[Seedream] Trying size={size_label} model={IMG_MODEL}', file=sys.stderr)
            print(f'[Seedream] Endpoint={IMG_URL}', file=sys.stderr)
            raw = http_post(IMG_URL, headers, json.dumps(body).encode(), proxy_url, 180)

            resp = json.loads(raw)
            print(f'[Seedream] Response keys: {list(resp.keys())}', file=sys.stderr)

            if 'error' in resp:
                err = resp['error']
                print(f'[Seedream] ⚠️ API error (size={size_label}): code={err.get("code")} msg={err.get("message","")}', file=sys.stderr)
                # If it's a size/param error (400-class), try next size
                # If it's auth (401), give up immediately
                if err.get('code') in ('InvalidParameter', 'InvalidAPIKey', 'AuthenticationError'):
                    if 'Auth' in str(err.get('code','')) or 'APIKey' in str(err.get('code','')):
                        break  # auth failure — no point retrying
                last_err = str(err)
                continue

            items = resp.get('data', [])
            if not items:
                print(f'[Seedream] ⚠️ Empty data array (size={size_label})', file=sys.stderr)
                last_err = 'Empty data array'
                continue

            item = items[0]
            if item.get('url'):
                img_url = item['url']
                print(f'[Seedream] ✅ Got URL (size={size_label}), downloading...', file=sys.stderr)
                opener = make_opener(proxy_url)
                with opener.open(img_url, timeout=90) as r:
                    data = r.read()
                print(f'[Seedream] ✅ Downloaded {len(data)//1024}KB (size={size_label})', file=sys.stderr)
                return data
            elif item.get('b64_json'):
                print(f'[Seedream] ✅ Got b64_json (size={size_label})', file=sys.stderr)
                return base64.b64decode(item['b64_json'])
            else:
                print(f'[Seedream] ⚠️ No url or b64_json in item: {list(item.keys())}', file=sys.stderr)
                last_err = f'No image data. Item keys: {list(item.keys())}'
                # Don't retry size — this is probably a content issue
                break

        except urllib.error.HTTPError as e:
            body_txt = ''
            try: body_txt = e.read().decode('utf-8', errors='replace')[:400]
            except: pass
            print(f'[Seedream] ❌ HTTP {e.code} (size={size_label}): {body_txt[:200]}', file=sys.stderr)
            last_err = f'HTTP {e.code}: {body_txt}'
            if e.code in (401, 403):
                print('[Seedream] ❌ Auth error — check API key in Settings', file=sys.stderr)
                break   # auth failure: no retrying
            # 400/422 = bad params → try next size candidate
            continue
        except Exception as e:
            print(f'[Seedream] ❌ Exception (size={size_label}): {type(e).__name__}: {e}', file=sys.stderr)
            last_err = f'{type(e).__name__}: {e}'
            continue

    print(f'[Seedream] ❌ All candidates failed. Last error: {last_err}', file=sys.stderr)
    return None


# ═══════════════════════════════════════════════════════════════════════════════
# Drawing helpers
# ═══════════════════════════════════════════════════════════════════════════════

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=None):
    color = color or C_BG
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_bg_image(slide, img_path, alpha_overlay=0.55):
    """Add a full-bleed background image to slide with a dark overlay for text readability."""
    if not img_path or not os.path.exists(img_path):
        return
    try:
        # Add image as background
        pic = slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)
        # Move to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        # Add semi-transparent dark overlay for text readability
        overlay = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        r_val = int(0x11 + (0xFF - 0x11) * 0)
        overlay.fill.fore_color.rgb = RGBColor(0x09, 0x0E, 0x1A)
        overlay.line.fill.background()
        _set_alpha(overlay, int((1.0 - alpha_overlay) * 100))
    except Exception as e:
        print(f'[BgImage] {e}', file=sys.stderr)


def rect(slide, x, y, w, h, fill_color, line_color=None, line_w=None, rounded=False):
    sid = 5 if rounded else 1
    s = slide.shapes.add_shape(sid, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color: s.line.color.rgb = line_color; s.line.width = line_w or Pt(1)
    else: s.line.fill.background()
    return s

def _set_line_spacing(paragraph, spacing=1.4):
    """Set line spacing on a paragraph via XML (pptx library workaround)."""
    try:
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _et
        pPr = paragraph._p.get_or_add_pPr()
        lnSpc = pPr.find(_qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = _et.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = lnSpc.find(_qn('a:spcPct'))
        if spcPct is None:
            spcPct = _et.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))
    except Exception:
        pass


def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True,
        italic=False, line_spacing=1.35):
    """Add a text box. Always word-wraps and sets consistent line spacing."""
    if not text: return None
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True   # ALWAYS on — prevents overflow off slide edge
    tf.margin_left   = Pt(0)
    tf.margin_right  = Pt(0)
    tf.margin_top    = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    run.font.name  = 'Arial'   # universal font — no substitution artifacts
    _set_line_spacing(p, line_spacing)
    return tb

def txt_multi(slide, lines, x, y, w, h, size, color, bold=False, line_spacing=1.4):
    """Multiple lines in one text frame, each on its own paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = str(line)
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.name  = 'Arial'
        _set_line_spacing(p, line_spacing)
    return tb

def connector(slide, x, y, w, color):
    """Horizontal separator line."""
    s = slide.shapes.add_shape(1, x, y, w, Pt(1))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def icon_circle(slide, x, y, size, color, alpha=20):
    """Small filled circle for icon background."""
    s = slide.shapes.add_shape(9, x, y, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    _set_alpha(s, alpha)
    return s

def _set_alpha(shape, alpha_pct):
    try:
        sf = shape.fill._xPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr')) or sf.find(qn('a:prstClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha')); a.set('val', str(int(alpha_pct*1000)))
    except: pass

def add_radial_circle(slide, x, y, w, h, color_from='3B82F6', color_to='111827', alpha_from=70):
    """Radial gradient decorative circle (cover / end slides)."""
    shape = slide.shapes.add_shape(9, x, y, w, h)
    shape.line.fill.background()
    spPr = shape._element.spPr
    for tag in ['a:solidFill','a:noFill','a:gradFill','a:pattFill']:
        for el in spPr.findall(qn(tag)): spPr.remove(el)
    gf = etree.SubElement(spPr, qn('a:gradFill'))
    gf.set('rotWithShape','1')
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    gs0 = etree.SubElement(gsLst, qn('a:gs')); gs0.set('pos','0')
    c0 = etree.SubElement(gs0, qn('a:srgbClr')); c0.set('val', color_from)
    al = etree.SubElement(c0, qn('a:alpha')); al.set('val', str(alpha_from*1000))
    gs1 = etree.SubElement(gsLst, qn('a:gs')); gs1.set('pos','100000')
    c1 = etree.SubElement(gs1, qn('a:srgbClr')); c1.set('val', color_to)
    al2 = etree.SubElement(c1, qn('a:alpha')); al2.set('val','0')
    path = etree.SubElement(gf, qn('a:path')); path.set('path','circle')
    ftr = etree.SubElement(path, qn('a:fillToRect'))
    ftr.set('l','50000'); ftr.set('t','50000'); ftr.set('r','50000'); ftr.set('b','50000')
    return shape

def slide_title(slide, text, y=Inches(0.83), size=36):
    """Standard slide title with consistent bottom padding."""
    txt(slide, text, Inches(0.83), y, SLIDE_W-Inches(1.66), Inches(0.78),
        size, C_WHITE, bold=True)
    # Thin accent line under title
    connector(slide, Inches(0.83), y+Inches(0.83), Inches(0.56), C_BLUE)


def slide_footer(slide, slide_num, total, label=''):
    """Thin footer bar at very bottom of every content slide."""
    W, H = SLIDE_W, SLIDE_H
    bar = rect(slide, 0, H-Inches(0.28), W, Inches(0.28), C_CARD)
    _set_alpha(bar, 60)
    # Left: report label
    txt(slide, '社媒竞对AI分析报告  ·  Powered by Doubao',
        Inches(0.42), H-Inches(0.26), Inches(6), Inches(0.24),
        8, C_T4, line_spacing=1.0)
    # Right: page number
    pg = f'{slide_num} / {total}'
    txt(slide, pg, W-Inches(1.5), H-Inches(0.26), Inches(1.3), Inches(0.24),
        9, C_T3, align=PP_ALIGN.RIGHT, line_spacing=1.0)

def mini_icon(slide, x, y, color):
    """Small 0.33" square icon circle badge."""
    s = slide.shapes.add_shape(1, x, y, Pt(22), Pt(22))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE RENDERERS — Doubao Reference Style
# ═══════════════════════════════════════════════════════════════════════════════

def make_cover(prs, s, date_str):
    """Cover: background image (if available) + radial gradient circles + title + brand teasers."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    # Background image with dark overlay (60% image visible)
    add_bg_image(slide, s.get('_img'), alpha_overlay=0.40)

    # Decorative gradient circles (top-left & bottom-right, subtle when image present)
    add_radial_circle(slide, Inches(-1.39), Inches(-1.39), Inches(5.56), Inches(5.56), '3B82F6', '111827', 80)
    add_radial_circle(slide, Inches(8.33),  Inches(2.78),  Inches(6.94), Inches(6.94), '2563EB', '111827', 65)

    heading    = s.get('heading','分析报告')
    subheading = s.get('subheading','')
    brands     = s.get('brands', [])       # list of {name, role}
    kpi_badges = s.get('kpi_badges', [])

    # Main title — 40pt, 1.2" height for up to 2 lines
    txt(slide, heading, Inches(1.11), Inches(2.0), Inches(11.11), Inches(1.25),
        40, C_WHITE, bold=True, line_spacing=1.2)

    # Horizontal separator line
    connector(slide, Inches(5.97), Inches(3.26), Inches(1.39), RGBColor(0x4B,0x55,0x63))

    # Subtitle
    if subheading:
        txt(slide, subheading, Inches(1.11), Inches(3.4), Inches(11.11), Inches(0.65),
            17, C_T2, line_spacing=1.3)

    # Footer line + date + KPI badges
    connector(slide, Inches(1.11), Inches(6.11), Inches(11.11), RGBColor(0x37,0x41,0x51))
    date_label = date_str if date_str else '社媒竞对监控 · AI智能分析报告'
    txt(slide, date_label, Inches(1.11), Inches(6.22), Inches(5.0), Inches(0.35),
        11, C_T4, line_spacing=1.0)
    # KPI badges on right side
    if kpi_badges:
        badges_txt = '  ·  '.join(str(b) for b in kpi_badges[:4])
        txt(slide, badges_txt, Inches(6.5), Inches(6.22), Inches(6.0), Inches(0.35),
            11, C_BLUE_M, align=PP_ALIGN.RIGHT, line_spacing=1.0)

    # Brand teaser cards at bottom (if brands provided)
    bx = Inches(1.39)
    for i, brand in enumerate(brands[:3]):
        bw = Inches(3.33)
        card_top = Inches(4.31)
        card_h   = Inches(1.53)
        card = rect(slide, bx, card_top, bw, card_h, C_CARD)
        # Colored left accent bar
        rect(slide, bx, card_top, Pt(4), card_h, BRAND_COLORS[i][0])
        # Mini colored icon circle
        ic_size = Pt(24)
        ic = slide.shapes.add_shape(9, bx+Inches(0.21), card_top+Inches(0.18), ic_size, ic_size)
        ic.fill.solid(); ic.fill.fore_color.rgb = BRAND_COLORS[i][0]; ic.line.fill.background()
        _set_alpha(ic, 100)
        # Brand name (bold, 15pt)
        name = brand.get('name','') if isinstance(brand,dict) else str(brand)
        role = brand.get('role','') if isinstance(brand,dict) else ''
        txt(slide, name, bx+Inches(0.6), card_top+Inches(0.14), bw-Inches(0.75), Inches(0.45),
            15, C_T1, bold=True, line_spacing=1.1)
        # Separator
        connector(slide, bx+Inches(0.21), card_top+Inches(0.65), bw-Inches(0.35),
                  RGBColor(0x37,0x41,0x51))
        # Role tag (colored label, 12pt)
        if role:
            txt(slide, role, bx+Inches(0.21), card_top+Inches(0.72), bw-Inches(0.3),
                Inches(0.56), 12, BRAND_COLORS[i][1], line_spacing=1.3)
        bx += bw + Inches(0.28)


def make_agenda(prs, s, slide_num, total):
    """Agenda: numbered rows with blue badge (Slide 2 pattern)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','汇报目录'))

    items = s.get('items', [])
    max_items = min(len(items), 7)
    # Distribute rows across available space
    avail_y = H - Inches(1.94) - Inches(0.32)
    row_h   = avail_y / max(max_items, 1)
    row_h   = max(min(row_h, Inches(0.85)), Inches(0.56))
    y0 = Inches(1.94)

    for i, item in enumerate(items[:7]):
        y = y0 + i * row_h
        if y + row_h > H - Inches(0.32): break
        gap = Pt(3)
        # Alternating subtle row bg
        row = rect(slide, Inches(0.83), y+gap/2, Inches(11.67), row_h-gap, C_CARD)
        _set_alpha(row, i%2==0 and 40 or 25)
        # Colored left accent bar
        rect(slide, Inches(0.83), y+gap/2, Pt(4), row_h-gap, C_BLUE)
        # Number badge (blue filled square)
        badge_size = min(row_h*0.55, Inches(0.39))
        bx = Inches(1.08); by = y + (row_h - badge_size)/2
        ic = rect(slide, bx, by, badge_size, badge_size, C_BLUE)
        # Number text (11pt white bold, fitted to badge)
        txt(slide, f'{i+1:02d}', bx, by, badge_size, badge_size,
            11, C_WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Item text — tall enough for wrapping
        txt(slide, item, Inches(1.72), y+Pt(4), Inches(9.97), row_h-Pt(8),
            15, C_T1, line_spacing=1.25)

    slide_footer(slide, slide_num, total)


def make_section(prs, s, slide_num, total):
    """Section divider: full background image (85% visible) + large chapter label."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    # Full-bleed background image (85% visible) or gradient circle fallback
    if s.get('_img') and os.path.exists(s.get('_img','')):
        add_bg_image(slide, s['_img'], alpha_overlay=0.15)
    else:
        add_radial_circle(slide, Inches(-0.5), Inches(-0.5), Inches(4.0), Inches(4.0), '3B82F6', '111827', 40)

    label = s.get('label', f'PART {slide_num:02d}')
    # Large label
    txt(slide, label, Inches(0.83), Inches(1.94), Inches(8), Inches(1.67),
        72, C_BLUE, bold=True)
    # Thin separator
    connector(slide, Inches(0.83), Inches(3.61), Inches(2.78), C_BLUE)
    # Main heading
    txt(slide, s.get('heading',''), Inches(0.83), Inches(3.78),
        Inches(11.67), Inches(2.5), 36, C_WHITE, bold=True)
    slide_footer(slide, slide_num, total)


def make_profile_cards(prs, s, slide_num, total):
    """3 white entity cards with colored header bars (Slide 4 pattern).
    Each card: white bg, colored top bar, brand name, then subfields with colored labels."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','品牌定位分析'), y=Inches(0.56), size=36)

    entities = s.get('entities', [])[:3]
    n = len(entities)
    if n == 0: return

    card_w = Inches(3.61)
    card_h = H - Inches(1.53) - Inches(0.28)
    bar_h  = Inches(0.69)
    starts = [Inches(0.83), Inches(4.86), Inches(8.89)]

    for i, entity in enumerate(entities):
        x = starts[i]
        color_main, color_light, color_mid = BRAND_COLORS[i]

        # White card background
        card = rect(slide, x, Inches(1.53), card_w, card_h, C_WHITE)
        _set_alpha(card, 5)  # subtle white card on dark bg

        # Colored header bar
        hdr_bar = rect(slide, x, Inches(1.53), card_w, bar_h, color_main)

        # Brand icon circle in header
        ic = slide.shapes.add_shape(9, x+Inches(0.21), Inches(1.67), Pt(28), Pt(28))
        ic.fill.solid(); ic.fill.fore_color.rgb = C_WHITE; ic.line.fill.background()
        _set_alpha(ic, 30)

        # Brand name in header
        name = entity.get('name','')
        role_tag = entity.get('tag','')
        hdr_label = f'{name}  ({role_tag})' if role_tag else name
        txt(slide, hdr_label, x+Inches(0.64), Inches(1.64), card_w-Inches(0.75), Inches(0.47),
            15, C_WHITE, bold=True)

        # Sub-fields — each field gets a label + value pair with auto-wrap
        fields = entity.get('fields', [])  # list of {label, value}
        # Compute dynamic row height: distribute available card body space evenly
        card_body_top = Inches(2.43)
        card_body_h   = card_h - (card_body_top - Inches(1.53)) - Inches(0.15)
        n_fields = min(len(fields), 5)
        if n_fields == 0:
            n_fields = 1
        field_h = card_body_h / n_fields
        # Cap: no taller than 1.1", no shorter than 0.7"
        field_h = max(min(field_h, Inches(1.1)), Inches(0.68))

        fy = card_body_top
        for j, field in enumerate(fields[:5]):
            fl = field.get('label','') if isinstance(field,dict) else ''
            fv = field.get('value','') if isinstance(field,dict) else str(field)
            # Icon dot
            ic2 = rect(slide, x+Inches(0.21), fy+Inches(0.08), Pt(14), Pt(14), color_main)
            # Label row
            txt(slide, fl, x+Inches(0.46), fy,
                card_w-Inches(0.6), Inches(0.30), 11, color_light, bold=True, line_spacing=1.1)
            # Value — taller box, auto-wrapping 11pt text
            val_h = field_h - Inches(0.32)
            txt(slide, fv, x+Inches(0.46), fy+Inches(0.28),
                card_w-Inches(0.6), val_h, 11, C_T2, line_spacing=1.3)
            fy += field_h
            if fy + field_h > Inches(1.53) + card_h: break

    slide_footer(slide, slide_num, total)


def make_content_strategy(prs, s, slide_num, total):
    """3 dark cards with image placeholder + colored strategy label (Slide 5 pattern)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','内容策略对比'), y=Inches(0.69), size=32)

    entities = s.get('entities', [])[:3]
    card_w  = Inches(3.61)
    card_y  = Inches(1.81)
    card_h  = H - card_y - Inches(0.28)
    img_h   = Inches(1.94)
    starts  = [Inches(0.83), Inches(4.86), Inches(8.89)]
    img_map = s.get('_img_map', {})

    for i, entity in enumerate(entities):
        x = starts[i]
        color_main, color_light, color_mid = BRAND_COLORS[i]

        # Dark card bg
        rect(slide, x, card_y, card_w, card_h, C_CARD2)

        # Brand header row
        ic = rect(slide, x+Inches(0.28), card_y+Inches(0.22), Pt(22), Pt(22), color_main)
        name = entity.get('name','')
        txt(slide, name, x+Inches(0.67), card_y+Inches(0.2), card_w-Inches(0.85), Inches(0.42),
            18, C_WHITE, bold=True)

        # Separator line
        connector(slide, x+Inches(0.28), card_y+Inches(0.65), card_w-Inches(0.56), color_main)

        # Image area (or colored placeholder)
        img_path = img_map.get(i)
        img_y = card_y+Inches(0.7)
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, x+Inches(0.28), img_y,
                                         card_w-Inches(0.56), img_h)
            except: pass
        else:
            # Colored image placeholder rectangle
            ph = rect(slide, x+Inches(0.28), img_y, card_w-Inches(0.56), img_h, C_CARD3)

        # Strategy label (colored)
        label = entity.get('strategy_label','')
        label_y = img_y + img_h + Inches(0.10)
        if label:
            txt(slide, label, x+Inches(0.28), label_y,
                card_w-Inches(0.42), Inches(0.36), 13, color_mid, bold=True, line_spacing=1.1)

        # Bullet points — dynamic height sharing remaining card space
        points  = entity.get('points', [])
        remain  = card_y + card_h - label_y - Inches(0.42) - Inches(0.08)
        n_pts   = min(len(points), 3)
        pt_h    = (remain / n_pts) if n_pts > 0 else Inches(0.7)
        pt_h    = max(min(pt_h, Inches(1.0)), Inches(0.55))
        py      = label_y + Inches(0.42)
        for pt in points[:3]:
            if py + pt_h > card_y + card_h: break
            # Colored left bar accent
            rect(slide, x+Inches(0.28), py+Pt(3), Pt(3), pt_h-Pt(6), color_main)
            txt(slide, pt.strip(), x+Inches(0.46), py,
                card_w-Inches(0.6), pt_h, 11, C_T2, line_spacing=1.3)
            py += pt_h

    slide_footer(slide, slide_num, total)


def make_dual_panel(prs, s, slide_num, total):
    """50/50 split with two content panels. Image used as right-panel background if available."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading',''), y=Inches(0.56), size=36)

    panels = s.get('panels', [])
    if not panels:
        # Fall back to col_a / col_b
        panels = [
            {'title': s.get('col_a_title','左栏'), 'items': s.get('col_a',[]),
             'type': 'bullets'},
            {'title': s.get('col_b_title','右栏'), 'items': s.get('col_b',[]),
             'type': 'chain'},
        ]

    panel_w  = (W - Inches(0.5)) / 2 - Inches(0.05)
    panel_y  = Inches(1.53)
    panel_h  = H - panel_y - Inches(0.28)
    starts   = [Inches(0.25), Inches(0.25) + panel_w + Inches(0.1)]

    for i, panel in enumerate(panels[:2]):
        x = starts[i]
        # White card (low alpha)
        card = rect(slide, x, panel_y, panel_w, panel_h, C_WHITE)
        _set_alpha(card, 5)

        # Panel header icon
        ic = rect(slide, x+Inches(0.28), panel_y+Inches(0.28), Pt(22), Pt(22), C_BLUE)
        txt(slide, panel.get('title',''), x+Inches(0.67), panel_y+Inches(0.25),
            panel_w-Inches(0.85), Inches(0.42), 18, C_T1, bold=True)

        # Separator
        connector(slide, x+Inches(0.28), panel_y+Inches(0.76), panel_w-Inches(0.56),
                  RGBColor(0x37,0x41,0x51))

        ptype = panel.get('type','bullets')
        items = panel.get('items',[])

        if ptype == 'chain':
            # Conversion chain style (with → separators)
            cy = panel_y + Inches(1.0)
            for j, item in enumerate(items[:3]):
                color = BRAND_COLORS[j][0]
                row = rect(slide, x+Inches(0.28), cy, panel_w-Inches(0.56), Inches(1.25), color)
                _set_alpha(row, 15)
                # Colored left stripe
                rect(slide, x+Inches(0.28), cy, Pt(4), Inches(1.25), color)
                # Icon
                ic2 = rect(slide, x+Inches(0.46), cy+Inches(0.38), Pt(22), Pt(22), color)
                title_txt = item.get('title','') if isinstance(item,dict) else str(item)
                desc_txt  = item.get('desc','')  if isinstance(item,dict) else ''
                txt(slide, title_txt, x+Inches(0.86), cy+Inches(0.1),
                    panel_w-Inches(1.1), Inches(0.39), 16, color, bold=True)
                if desc_txt:
                    txt(slide, desc_txt, x+Inches(0.46), cy+Inches(0.52),
                        panel_w-Inches(0.7), Inches(0.56), 13, C_T2)
                cy += Inches(1.39)
        else:
            # Bullets style — dynamic height
            bullet_items = items[:6]
            panel_body_start = panel_y + Inches(0.95)
            panel_body_end   = panel_y + panel_h - Inches(0.1)
            avail_h = panel_body_end - panel_body_start
            n_b = len(bullet_items)
            row_h = (avail_h / n_b) if n_b > 0 else Inches(0.75)
            row_h = max(min(row_h, Inches(1.05)), Inches(0.50))
            py = panel_body_start
            for j, item in enumerate(bullet_items):
                if py + row_h > panel_body_end + Inches(0.1): break
                label = item.get('label','') if isinstance(item,dict) else ''
                desc  = item.get('desc','')  if isinstance(item,dict) else str(item)
                if label:
                    lh = Inches(0.26)
                    txt(slide, label, x+Inches(0.28), py, panel_w-Inches(0.38),
                        lh, 12, C_T1, bold=True, line_spacing=1.1)
                    txt(slide, desc, x+Inches(0.28), py+lh, panel_w-Inches(0.38),
                        row_h-lh, 11, C_T3, line_spacing=1.3)
                else:
                    dot = rect(slide, x+Inches(0.28), py+Inches(0.06), Pt(5), Pt(5), C_BLUE)
                    txt(slide, desc, x+Inches(0.46), py, panel_w-Inches(0.54),
                        row_h, 11, C_T2, line_spacing=1.3)
                py += row_h

    slide_footer(slide, slide_num, total)


def make_flow_insight(prs, s, slide_num, total):
    """Wide banner highlight + 3 colored insight cards below. Optional right-panel image."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    # Optional: add background image as subtle texture (30% visible)
    if s.get('_img') and os.path.exists(s.get('_img','')):
        add_bg_image(slide, s['_img'], alpha_overlay=0.70)

    slide_title(slide, s.get('heading','市场洞察'), y=Inches(0.83), size=36)

    # Top wide banner (evolution flow or key insight)
    top_banner = s.get('banner', {})
    if top_banner:
        banner = rect(slide, Inches(0.83), Inches(1.94), W-Inches(1.66), Inches(2.08), C_WHITE)
        _set_alpha(banner, 5)
        # Banner content: flow items or insight text
        flow_items = top_banner.get('flow_items', [])
        if flow_items:
            item_w = (W - Inches(1.66)) / len(flow_items)
            for j, fi in enumerate(flow_items[:4]):
                ix = Inches(0.83) + j * item_w
                # Icon
                ic_s = slide.shapes.add_shape(9, ix+Inches(0.28), Inches(2.22), Inches(0.67), Inches(0.67))
                ic_s.fill.solid(); ic_s.fill.fore_color.rgb = C_BLUE; ic_s.line.fill.background()
                _set_alpha(ic_s, 20)
                # Number
                txt(slide, fi.get('num',f'{j+1:02d}'), ix+Inches(0.28), Inches(2.99),
                    Inches(2.22), Inches(0.83), 16, C_T1, bold=True)
                txt(slide, fi.get('desc',''), ix+Inches(0.28), Inches(3.22),
                    Inches(2.22), Inches(0.56), 12, C_T3)
                # Arrow between items
                if j < len(flow_items)-1:
                    txt(slide, '→', ix+item_w-Inches(0.3), Inches(2.38),
                        Inches(0.4), Inches(0.33), 14, C_BLUE)
        else:
            insight_text = top_banner.get('text','')
            if insight_text:
                txt(slide, insight_text, Inches(1.11), Inches(2.08),
                    W-Inches(2.22), Inches(1.8), 14, C_T2)
    else:
        # Push cards up
        pass

    # Bottom 3 colored insight cards
    cards = s.get('cards', [])[:3]
    card_w = (W - Inches(0.5)) / 3 - Inches(0.08)
    card_y = Inches(4.31)
    card_h = H - card_y - Inches(0.28)
    xs = [Inches(0.25) + i*(card_w+Inches(0.08)) for i in range(3)]

    for j, card in enumerate(cards[:3]):
        x = xs[j]
        color = BRAND_COLORS[j][0]
        card_bg = rect(slide, x, card_y, card_w, card_h, color)
        _set_alpha(card_bg, 15)
        rect(slide, x, card_y, Pt(4), card_h, color)

        # Card header icon
        ic = rect(slide, x+Inches(0.28), card_y+Inches(0.28), Pt(20), Pt(20), color)
        card_title = card.get('title','') if isinstance(card,dict) else ''
        txt(slide, card_title, x+Inches(0.62), card_y+Inches(0.22),
            card_w-Inches(0.78), Inches(0.46), 16, C_T1, bold=True, line_spacing=1.2)

        connector(slide, x+Inches(0.28), card_y+Inches(0.76), card_w-Inches(0.42),
                  RGBColor(0x37,0x41,0x51))

        body = card.get('body','') if isinstance(card,dict) else str(card)
        # Body gets the full remaining height with generous wrapping
        body_h = card_h - Inches(0.88)
        txt(slide, body, x+Inches(0.28), card_y+Inches(0.88),
            card_w-Inches(0.42), body_h, 12, C_T2, line_spacing=1.4)

    slide_footer(slide, slide_num, total)


def make_three_analysis(prs, s, slide_num, total):
    """3 white analysis cards (Slides 9/11/12 pattern): icon+title+separator+body."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading',''), y=Inches(0.69), size=36)

    cards = s.get('cards', [])[:3]
    card_w = (W - Inches(0.5)) / 3 - Inches(0.08)
    card_y = Inches(1.81)
    card_h = H - card_y - Inches(0.28)
    xs = [Inches(0.25) + i*(card_w+Inches(0.08)) for i in range(3)]

    for j, card in enumerate(cards[:3]):
        x = xs[j]
        color = BRAND_COLORS[j][0]
        color_l = BRAND_COLORS[j][1]

        # Card background (subtle white)
        cbg = rect(slide, x, card_y, card_w, card_h, C_WHITE)
        _set_alpha(cbg, 5)

        # Top icon + title
        ic = rect(slide, x+Inches(0.28), card_y+Inches(0.28), Pt(22), Pt(22), color)
        card_title = card.get('title','') if isinstance(card,dict) else ''
        txt(slide, card_title, x+Inches(0.67), card_y+Inches(0.25),
            card_w-Inches(0.85), Inches(0.42), 18, C_T1, bold=True)

        # Separator
        connector(slide, x+Inches(0.28), card_y+Inches(0.76), card_w-Inches(0.56),
                  RGBColor(0x37,0x41,0x51))

        # Sub-items with colored labels
        items = card.get('items', [])
        if items:
            # Dynamic row height: share card body between items
            card_body_start = card_y + Inches(0.95)
            card_body_end   = card_y + card_h - Inches(0.1)
            avail_h  = card_body_end - card_body_start
            n_items  = min(len(items), 5)
            row_h    = avail_h / n_items if n_items > 0 else Inches(0.9)
            row_h    = max(min(row_h, Inches(1.15)), Inches(0.62))
            iy = card_body_start
            for k, item in enumerate(items[:5]):
                if iy + row_h > card_body_end + Inches(0.1): break
                ilabel = item.get('label','') if isinstance(item,dict) else ''
                ibody  = item.get('body','')  if isinstance(item,dict) else str(item)
                if ilabel:
                    lh = Inches(0.26)
                    txt(slide, ilabel, x+Inches(0.28), iy,
                        card_w-Inches(0.4), lh, 12, color_l, bold=True, line_spacing=1.1)
                    txt(slide, ibody, x+Inches(0.28), iy+lh,
                        card_w-Inches(0.4), row_h-lh, 11, C_T3, line_spacing=1.3)
                else:
                    dot = rect(slide, x+Inches(0.28), iy+Inches(0.08), Pt(5), Pt(5), color)
                    txt(slide, ibody, x+Inches(0.48), iy,
                        card_w-Inches(0.58), row_h, 11, C_T2, line_spacing=1.3)
                iy += row_h

        else:
            body = card.get('body','') if isinstance(card,dict) else str(card)
            txt(slide, body, x+Inches(0.28), card_y+Inches(0.9),
                card_w-Inches(0.4), card_h-Inches(1.05), 12, C_T2, line_spacing=1.35)

    slide_footer(slide, slide_num, total)


def make_comparison_table(prs, s, slide_num, total):
    """Blue header + comparison table rows (Slide 7 pattern)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading',''), y=Inches(0.83), size=36)

    rows = s.get('rows', [])
    cols = s.get('cols', [])   # column headers
    if not rows and not cols: return

    table_y = Inches(1.81)
    table_h = H - table_y - Inches(0.28)
    n_rows  = len(rows) + 1   # +1 for header
    n_cols  = max(len(cols), max((len(r) for r in rows), default=0))
    if n_cols < 2: n_cols = 2

    col_w   = (W - Inches(1.66)) / n_cols
    row_h   = table_h / n_rows

    # Header row
    rect(slide, Inches(0.83), table_y, W-Inches(1.66), row_h, C_BLUE)
    for j, col in enumerate(cols[:n_cols]):
        cx = Inches(0.83) + j * col_w
        txt(slide, col, cx+Inches(0.1), table_y+Inches(0.1),
            col_w-Inches(0.2), row_h-Inches(0.1), 14, C_WHITE, bold=True)

    # Data rows
    for i, row in enumerate(rows[:8]):
        ry = table_y + (i+1) * row_h
        row_bg = rect(slide, Inches(0.83), ry, W-Inches(1.66), row_h,
                      C_CARD if i%2==0 else C_CARD2)
        for j, cell in enumerate(row[:n_cols]):
            cx = Inches(0.83) + j * col_w
            cell_color = BRAND_COLORS[j][1] if j > 0 else C_T1
            # Left column: bold dimension label; other cols: normal content
            is_dim = (j == 0)
            txt(slide, str(cell), cx+Pt(8), ry+Pt(5),
                col_w-Pt(12), row_h-Pt(8), 11 if not is_dim else 12,
                cell_color, bold=is_dim, line_spacing=1.3)

    slide_footer(slide, slide_num, total)


def make_metrics(prs, s, slide_num, total):
    """Large KPI number cards (data dashboard)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','数据看板'), y=Inches(0.56))
    if s.get('subtitle'):
        txt(slide, s['subtitle'], Inches(0.83), Inches(1.35), W-Inches(1.66), Inches(0.35),
            12, C_T4)

    metrics = s.get('metrics', [])[:4]
    n = len(metrics)
    if n == 0: return
    card_w = (W - Inches(0.5)) / n - Inches(0.1)
    card_y = Inches(1.6)
    card_h = H - card_y - Inches(0.28)

    for i, m in enumerate(metrics):
        x = Inches(0.25) + i*(card_w+Inches(0.1))
        color = BRAND_COLORS[i % len(BRAND_COLORS)][0]
        color_l = BRAND_COLORS[i % len(BRAND_COLORS)][1]

        card = rect(slide, x, card_y, card_w, card_h, C_CARD)
        rect(slide, x, card_y, card_w, Pt(4), color)

        # Decorative circle
        ic = slide.shapes.add_shape(9, x+card_w-Pt(38), card_y+Pt(10), Pt(28), Pt(28))
        ic.fill.solid(); ic.fill.fore_color.rgb = color; ic.line.fill.background()
        _set_alpha(ic, 20)

        # Big value
        txt(slide, str(m.get('value','—')), x+Pt(16), card_y+Pt(16),
            card_w-Pt(20), Pt(78), 38, color, bold=True, line_spacing=1.0)
        # Separator line
        connector(slide, x+Pt(12), card_y+card_h*0.52, card_w-Pt(18),
                  RGBColor(0x37,0x41,0x51))
        txt(slide, str(m.get('label','')), x+Pt(16), card_y+card_h*0.55,
            card_w-Pt(20), Pt(34), 13, C_T1, line_spacing=1.2)
        if m.get('trend'):
            txt(slide, str(m['trend']), x+Pt(16), card_y+card_h*0.73,
                card_w-Pt(20), Pt(32), 11, C_GREEN_M, line_spacing=1.2)

    slide_footer(slide, slide_num, total)


def make_bar_chart(prs, s, slide_num, total):
    """Horizontal bar chart with value labels."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','数据对比'), y=Inches(0.56))
    if s.get('subtitle'):
        txt(slide, s['subtitle'], Inches(0.83), Inches(1.35), W-Inches(1.66), Inches(0.28),
            11, C_T4)

    bars = s.get('bars', [])
    if not bars: return

    chart_x = Inches(2.0)
    chart_w = W - Inches(2.8)
    chart_y = Inches(1.7)
    chart_h = H - chart_y - Inches(0.4)
    row_h   = chart_h / len(bars[:6])

    for i, bar in enumerate(bars[:6]):
        color   = BRAND_COLORS[i % len(BRAND_COLORS)][0]
        color_l = BRAND_COLORS[i % len(BRAND_COLORS)][1]
        y     = chart_y + i * row_h
        val   = float(bar.get('value', 0))
        maxv  = float(bar.get('max', 100))
        pct   = min(val/maxv, 1.0) if maxv > 0 else 0

        # Label — auto-wrap enabled, uses 14pt
        label = str(bar.get('label',''))
        txt(slide, label, Inches(0.12), y + row_h*0.08, Inches(1.78), row_h*0.84,
            13, C_T1, line_spacing=1.2)

        # Track (background bar)
        bth = max(row_h * 0.28, Pt(14))
        bty = y + (row_h - bth) / 2
        rect(slide, chart_x, bty, chart_w, bth, C_CARD2)

        # Fill bar
        fw = chart_w * pct
        if fw > Pt(4):
            fill = rect(slide, chart_x, bty, fw, bth, color)
            # Subtle gradient end cap
            endcap = rect(slide, chart_x+fw-Pt(4), bty, Pt(4), bth, color_l)

        # Value label — placed to the right of bar, or inside if bar is long enough
        val_str = f'{val:.0f}%' if maxv <= 100 else f'{val:.1f}'
        label_x = chart_x + fw + Pt(8) if pct < 0.85 else chart_x + fw - Inches(0.7)
        label_color = C_T1 if pct < 0.85 else C_WHITE
        txt(slide, val_str, label_x, bty, Inches(0.65), bth,
            12, label_color, bold=True, line_spacing=1.0)

    slide_footer(slide, slide_num, total)


def make_summary(prs, s, slide_num, total):
    """Recommendations: 3 dark cards + per-audience strategy (Slide 12 pattern)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    slide_title(slide, s.get('heading','核心启示与建议'), y=Inches(0.69), size=36)

    cards = s.get('cards', s.get('points', []))
    if not cards:
        # Fallback: use points as simple bullet list
        points = s.get('points', [])
        py = Inches(1.7)
        for i, pt in enumerate(points[:6]):
            color = BRAND_COLORS[i%len(BRAND_COLORS)][0]
            card = rect(slide, Inches(0.83), py, W-Inches(1.66), Inches(0.92), C_CARD3)
            rect(slide, Inches(0.83), py, Pt(4), Inches(0.92), color)
            nb = slide.shapes.add_shape(1, Inches(1.0), py+Inches(0.22), Pt(22), Pt(22))
            nb.fill.solid(); nb.fill.fore_color.rgb = color; nb.line.fill.background()
            txt(slide, str(i+1), Inches(1.0), py+Inches(0.19), Pt(22), Pt(24),
                10, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            txt(slide, str(pt).strip(), Inches(1.42), py+Inches(0.14),
                W-Inches(2.3), Inches(0.65), 13, C_T2, wrap=True)
            py += Inches(1.0)
        # CTA
        cta = s.get('cta','')
        if cta:
            rect(slide, 0, H-Inches(0.6), W, Inches(0.6), C_CARD)
            connector(slide, 0, H-Inches(0.6), W, C_BLUE)
            txt(slide, f'▶  {cta}', Inches(0.83), H-Inches(0.52), W-Inches(1.66), Inches(0.42),
                12, C_BLUE_M, bold=True)
        return

    slide_footer(slide, slide_num, total)

    # 3-card layout
    card_w = (W - Inches(0.5)) / min(len(cards),3) - Inches(0.08)
    card_y = Inches(1.81)
    card_h = H - card_y - Inches(0.3)

    for j, card in enumerate(cards[:3]):
        x = Inches(0.25) + j*(card_w+Inches(0.08))
        color = BRAND_COLORS[j][0]

        cbg = rect(slide, x, card_y, card_w, card_h, C_CARD3)
        # Icon + title
        ic = rect(slide, x+Inches(0.28), card_y+Inches(0.28), Pt(22), Pt(22), color)
        card_title = card.get('title','') if isinstance(card,dict) else str(card)
        txt(slide, card_title, x+Inches(0.67), card_y+Inches(0.25),
            card_w-Inches(0.85), Inches(0.49), 18, BRAND_COLORS[j][2], bold=True)
        connector(slide, x+Inches(0.28), card_y+Inches(0.8), card_w-Inches(0.56),
                  RGBColor(0x37,0x41,0x51))

        items = card.get('items',[]) if isinstance(card,dict) else []
        # Dynamic row height: share card body evenly
        body_start = card_y + Inches(0.95)
        body_end   = card_y + card_h - Inches(0.1)
        n_it = min(len(items), 4)
        row_h = ((body_end - body_start) / n_it) if n_it > 0 else Inches(1.0)
        row_h = max(min(row_h, Inches(1.2)), Inches(0.62))
        iy = body_start
        for k, item in enumerate(items[:4]):
            if iy + row_h > body_end + Inches(0.1): break
            ilabel = item.get('label','') if isinstance(item,dict) else ''
            ibody  = item.get('body','')  if isinstance(item,dict) else str(item)
            if ilabel:
                ic2 = slide.shapes.add_shape(9, x+Inches(0.28), iy+Inches(0.06),
                                              Pt(13), Pt(13))
                ic2.fill.solid(); ic2.fill.fore_color.rgb = color; ic2.line.fill.background()
                lh = Inches(0.26)
                txt(slide, ilabel, x+Inches(0.50), iy, card_w-Inches(0.62),
                    lh, 12, C_T1, bold=True, line_spacing=1.1)
                txt(slide, ibody, x+Inches(0.50), iy+lh, card_w-Inches(0.62),
                    row_h-lh, 11, C_T3, line_spacing=1.3)
            else:
                dot = rect(slide, x+Inches(0.28), iy+Inches(0.08), Pt(5), Pt(5), color)
                txt(slide, ibody, x+Inches(0.46), iy, card_w-Inches(0.56),
                    row_h, 11, C_T2, line_spacing=1.3)
            iy += row_h


def make_end(prs):
    """End slide: gradient circles + centered thank you + contact (Slide 13 pattern)."""
    slide = blank(prs); bg(slide)
    W, H = SLIDE_W, SLIDE_H

    add_radial_circle(slide, Inches(-1.39), Inches(-1.39), Inches(5.56), Inches(5.56), '3B82F6', '18181B', 80)
    add_radial_circle(slide, Inches(8.33),  Inches(2.78),  Inches(6.94), Inches(6.94), '2563EB', '18181B', 65)

    txt(slide, '感谢聆听', Inches(1.11), Inches(2.22), W-Inches(2.22), Inches(1.11),
        48, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, '欢迎交流探讨', Inches(1.11), Inches(3.33), W-Inches(2.22), Inches(0.56),
        20, C_T2, align=PP_ALIGN.CENTER)

    # Center line
    connector(slide, Inches(6.11), Inches(4.02), Inches(1.11), RGBColor(0x37,0x41,0x51))

    # Contact bar
    cbg = rect(slide, Inches(2.5), Inches(5.0), Inches(8.33), Inches(1.11), C_WHITE)
    _set_alpha(cbg, 6)
    txt(slide, 'contact@company.com', Inches(3.61), Inches(5.35), Inches(2.5), Inches(0.42),
        14, C_T1)
    txt(slide, '社媒竞对监控 · AI智能分析报告 · Powered by Doubao',
        Inches(1.11), Inches(5.5), W-Inches(2.22), Inches(0.4),
        11, C_T4, align=PP_ALIGN.CENTER)




# ═══════════════════════════════════════════════════════════════════════════════
# Step 1 – Strict JSON structure prompt + robust parsing + full validation
# ═══════════════════════════════════════════════════════════════════════════════

STRUCTURE_PROMPT = """你是专业商务PPT结构化AI。将社媒竞对分析报告100%转换为可直接生成PPT的纯JSON数据。

# 强制执行规则（每条都是硬性约束，违反则生成失败）
1. 仅输出纯JSON，第一个字符必须是 { 最后一个字符必须是 }，禁止任何markdown标记(```json)、解释文字、前后缀
0. 【严格禁止】禁止在任何JSON字段值中使用emoji字符（包括🔹✅🔸🚀✨⚡👉🔍📊等）；每页正文body/desc/value字段必须≥30字，禁止仅输出标题无正文
2. 所有文本字段必须填充真实有效内容，禁止空字符串""、禁止"待补充"/"暂无"/"示例"/"占位"等占位词
3. 100%基于原文提取，禁止编造数据；无精确数字时用原文定性描述（如"较高"/"领先"）代替
4. 单页只讲1个核心主题，禁止多主题合并；正文body/desc/value字段每条内容60-120字（中文），禁止少于30字的简短描述
5. entities/cards数组固定3个，panels数组固定2个，profile_cards的fields固定5个
6. 每张幻灯片必须填写 image_theme 字段：10字以内，精准描述本页核心视觉主题

# 幻灯片类型与字段规范（所有字段必填）

cover（封面）
  heading: 主标题≤25字
  subheading: 副标题，包含品牌名和分析范围≤45字
  brands: [{name:"品牌全名", role:"≤6字定位描述"}] 3个，来自原文实际品牌
  kpi_badges: ["分析X个账号","分析X条内容","X大维度"] 3-4个真实数字
  image_theme: "封面视觉主题"

agenda（目录）
  heading: "汇报目录"
  items: ["01. 章节名≤20字",...] 5-7条，必须覆盖报告所有章节
  image_theme: "目录页视觉"

section（章节分隔）
  heading: 章节主题≤20字
  label: "PART 01" 按序递增
  image_theme: "章节封面视觉"

profile_cards（品牌定位，3列白卡+彩色顶栏）
  heading: 页面标题
  entities: [{
    name: "品牌全名",
    tag: "≤6字定位标签",
    fields: [
      {label:"品牌位势",   value:"80-130字真实描述，必须包含具体内容，基于原文"},
      {label:"核心市场定位", value:"80-120字，品牌通过内容传递的差异化定位信息，基于原文"},
      {label:"核心目标市场", value:"80-120字，目标地理区域/行业垂类/企业规模，基于原文"},
      {label:"核心受众",   value:"80-120字，决策者职业/角色/需求特征，基于原文"},
      {label:"核心商业目标", value:"80-120字，内容驱动的转化路径与商业目标，基于原文"}
    ]
  }] 3个实体
  image_theme: "品牌定位对比"

content_strategy（内容策略对比，3列深色卡）
  heading: 页面标题
  entities: [{
    name: "品牌名",
    strategy_label: "核心打法标签≤12字",
    points: ["80-120字具体策略描述，基于原文实际内容，包含具体例子或数据", "第二条80-120字", "第三条80-120字"]
  }] 3个实体，points每个实体必须3条
  image_theme: "内容策略视觉"

dual_panel（双栏对比50/50）
  heading: 页面标题
  panels: [
    {title:"左栏标题≤12字", type:"bullets", items:[{label:"≤10字",desc:"60-100字具体描述，基于原文，包含可验证的具体内容"}] 3-4条},
    {title:"右栏标题≤12字", type:"chain",   items:[{title:"≤15字",desc:"60-100字具体描述，基于原文内容，不得笼统"}] 3条}
  ]
  image_theme: "双栏主题视觉"

flow_insight（流程洞察+底部3卡）
  heading: 页面标题
  banner: {flow_items:[{num:"01",desc:"25-40字演进描述，具体行动或趋势"}] 3-4步} 或 {text:"80-120字核心洞察"}
  cards: [{title:"≤12字",body:"80-120字具体洞察内容，来自原文实际分析，包含具体依据"}] 3张
  image_theme: "流程洞察视觉"

three_analysis（三栏分析卡）
  heading: 页面标题
  cards: [{
    title: "≤12字主题",
    items: [{label:"≤8字子标题", body:"60-100字具体分析内容，来自原文，禁止空泛"}] 4-5条
  }] 3张卡，每张items 4-5条
  image_theme: "分析主题视觉"

comparison_table（对比汇总表）
  heading: 页面标题
  cols: ["维度","品牌A名称","品牌B名称","品牌C名称"]
  rows: [["≤8字维度","A的30-50字真实描述，不得空泛","B的30-50字真实描述","C的30-50字真实描述"]] 7-9行
  image_theme: "对比表格视觉"

metrics（数据看板）
  heading: 标题
  subtitle: "≤30字数据来源说明"
  metrics: [{value:"真实数字或短词",label:"≤8字指标名",trend:"≤10字趋势"}] 3-4个
  image_theme: "数据看板视觉"

bar_chart（水平柱状图）
  heading: 标题
  subtitle: "≤25字数据说明"
  bars: [{label:"≤8字",value:数字,max:100,color_idx:0}] 4-6条
  image_theme: "数据对比视觉"

summary（总结建议，3列深色卡）
  heading: "核心结论与行动建议"
  cards: [{
    title: "≤10字主题",
    items: [{label:"≤8字要点", body:"80-120字可执行具体建议，来自原文分析结论，包含行动步骤"}] 3-4条
  }] 3张
  cta: "≤45字核心行动号召"
  image_theme: "总结建议视觉"

# 页面数量要求
cover×1, agenda×1, section≥2, profile_cards×1, content_strategy×1, comparison_table×1, flow_insight×1, summary×1
原文<500字: 总8-10页；500-2000字: 总12-15页；>2000字: 总15-18页

# 强制要求：正文内容字数下限
每个body/desc/value/points字段必须100-150字（中文），完整完句，绝对禁止在句子中途截断。
three_analysis的cards必须使用items数组格式（每张4-5个item），禁止使用简单的body字段替代。
summary必须使用cards数组格式（3张card，每张3-4个item），禁止使用points简单列表。

# 参考示例（严格对照此格式生成，字段名、结构、内容丰富度必须与示例保持一致）
以下是一套完整的高质量输出示例，你的输出必须达到相同的内容密度和结构完整性：
{"slides":[{"type":"cover","heading":"全球安防品牌领英运营策略竞对分析","subheading":"Hikvision / Dahua / Axis 三大品牌内容策略深度拆解","brands":[{"name":"Hikvision","role":"头部规模品牌"},{"name":"Dahua","role":"技术驱动品牌"},{"name":"Axis","role":"高端商用品牌"}],"kpi_badges":["分析3个账号","分析30条内容","6大维度对比","AI深度生成"],"image_theme":"城市安防科技夜景"},{"type":"agenda","heading":"汇报目录","items":["01. 品牌定位与核心受众画像","02. 内容策略与话题打法对比","03. 传播逻辑与发布节奏分析","04. 行业信号与竞争格局洞察","05. 数据指标横向对比看板","06. 核心结论与可执行建议"],"image_theme":"商务简约目录背景"},{"type":"section","heading":"品牌定位与核心受众","label":"PART 01","image_theme":"企业品牌识别视觉"},{"type":"profile_cards","heading":"三大品牌核心定位对比","entities":[{"name":"Hikvision","tag":"规模化领导者","fields":[{"label":"品牌位势","value":"全球出货量第一的视频监控品牌，以产品线宽度和成本优势构建护城河，在领英持续输出AIoT技术叙事，强调从摄像头到云端的完整解决方案能力，塑造行业技术标准制定者形象。"},{"label":"核心市场定位","value":"定位为全场景智能安防解决方案平台，覆盖政府、交通、零售、教育、能源等12个以上垂直行业，重点强调AI赋能的主动防御能力，区别于传统被动监控的叙事框架。"},{"label":"核心目标市场","value":"以亚太、中东、非洲为主要增长市场，深耕智慧城市和平安城市项目，着重进攻中型企业和政府集成商渠道，近期内容显示加速欧洲B2B市场渗透战略。"},{"label":"核心受众","value":"一线城市安防系统集成商（30-50岁技术决策层）、大型企业安全运营总监、政府智慧城市项目采购官员，以及持续寻找规模化部署成本解决方案的渠道合作伙伴。"},{"label":"核心商业目标","value":"通过高频内容输出强化品牌技术可信度，驱动集成商合作伙伴生态扩张，配合产品发布节奏制造行业话题声量，最终转化为全球经销商网络的持续询盘增长。"}]},{"name":"Dahua","tag":"技术深度品牌","fields":[{"label":"品牌位势","value":"全球第二大视频监控企业，以深度技术研发投入（营收15%用于研发）为核心壁垒，领英内容高度聚焦AI算法、边缘计算和智能分析技术，塑造技术驱动而非规模驱动的品牌认知。"},{"label":"核心市场定位","value":"定位为以AI视觉技术为核心的智慧物联网解决方案提供商，着重展示算法专利、技术认证和实验室级研发能力，强化与学术机构的技术共建叙事，区别于竞品的规模化叙事路径。"},{"label":"核心目标市场","value":"重点深耕中东、东南亚及拉美政府项目市场，同步强攻全球中高端零售和物流仓储行业，内容策略配合海外展会节奏集中释放，针对工程商渠道精准触达高价值项目。"},{"label":"核心受众","value":"系统集成商技术负责人、IT与安防融合项目的CTO级决策者、对AI算法落地有明确需求的行业垂直场景用户（零售损耗管理、仓储安全），以及寻求技术差异化的高端渠道商。"},{"label":"核心商业目标","value":"通过技术内容营销建立高端品牌溢价认知，吸引对价格不敏感的高价值项目客户，强化与系统集成商的技术绑定关系，推动Dahua智慧云平台SaaS化订阅模式的规模转化。"}]},{"name":"Axis","tag":"高端商用先锋","fields":[{"label":"品牌位势","value":"全球网络摄像头的发明者和高端商用安防标杆，依托施耐德电气母公司背书强化企业级品牌可信度，领英内容聚焦思想领导力和行业白皮书输出，塑造行业布道者而非产品推销者形象。"},{"label":"核心市场定位","value":"定位为企业级开放生态网络视频解决方案领导者，强调AXIS OS开放平台和第三方应用集成能力，重点突出零信任网络安全架构和隐私合规框架，契合欧美企业IT安全治理需求。"},{"label":"核心目标市场","value":"北美、西欧成熟市场的大型企业和关键基础设施保护，重点聚焦金融、医疗、教育和零售连锁行业，内容明显倾向合规性和安全认证话题，匹配欧盟GDPR和美国NDAA法规要求。"},{"label":"核心受众","value":"企业物理安全与IT安全融合决策者（CSO/CISO级别）、注重合规与隐私保护的北美欧洲企业安全经理、关注生态集成的方案商，以及对品牌溢价和长期售后体系有强需求的终端客户。"},{"label":"核心商业目标","value":"通过高质量思想领导力内容建立行业话语权，驱动高单价项目询盘转化，巩固AXIS Camera Station和AXIS Object Analytics等软件平台的用户粘性，推动从硬件销售向软件订阅模式迁移。"}]}],"image_theme":"企业品牌对比展示"},{"type":"content_strategy","heading":"内容打法与话题策略对比","entities":[{"name":"Hikvision","strategy_label":"规模输出+场景覆盖","points":["高频多话题并行策略：每周3-5条内容，话题横跨AIoT、智慧城市、零售AI等8个以上场景，以覆盖宽度换取细分受众触达，配合产品发布和展会节点集中爆发，构建持续的品牌声量底层。","以技术产品为主角的叙事框架：大量使用数据和参数强化技术可信度（如99.8%识别准确率），配合真实项目案例图片建立规模化部署的视觉证明，内容主角是产品而非客户故事。","渠道合作伙伴激活导向：相当比例的内容以集成商和渠道商为主要受众，强调培训、认证、联合销售资源，将领英作为B2B2C渠道激活工具而非纯粹的终端品牌建设平台。"]},{"name":"Dahua","strategy_label":"技术深度+垂直渗透","points":["技术深度内容驱动差异化认知：重点输出AI算法原理解析、技术白皮书摘要、实验室测试数据和专利技术图解，内容密度高、专业门槛高，有效筛选技术型决策者受众群体。","垂直行业场景深耕策略：围绕零售客流分析、仓储异常检测、交通流量管理等核心场景持续输出，每个场景配套痛点白皮书+解决方案视频+客户证言，形成完整的内容转化漏斗。","展会节点内容密集爆发：在Intersec、IFSEC等重要展会前后2周内容频率提升3倍，以预热-现场直播-成果回顾三段式覆盖完整展会传播周期，单次展会带来显著账号增粉。"]},{"name":"Axis","strategy_label":"思想领导力+生态运营","points":["思想领导力内容为核心武器：以行业趋势报告、政策解读（GDPR/NDAA合规）、技术观点文章为主，平均每篇内容字数超800字，大量引用第三方研究数据，建立高可信度行业媒体形象。","开放生态合作伙伴联合内容：定期发布与Milestone、Genetec等生态伙伴的联合解决方案内容，强调AXIS开放平台的集成能力，通过合作方背书扩大品牌可信度和触达受众范围。","隐私与安全合规内容占位：主动围绕数据隐私、网络安全合规、零信任架构等欧美客户高度关注议题持续输出，提前占据合规安防类的内容心智，强化产品差异化叙事护城河。"]}],"image_theme":"内容策略数字营销"},{"type":"section","heading":"传播逻辑与数据洞察","label":"PART 02","image_theme":"数据可视化科技感"},{"type":"comparison_table","heading":"三大品牌全景对比总览","cols":["对比维度","Hikvision","Dahua","Axis"],"rows":[["内容发布频率","每周3-5条，高频多点","每周2-3条，精准深度","每周1-2条，质量优先"],["核心话题方向","AIoT+智慧城市+展会","AI算法+垂直行业场景","合规+生态+思想领导力"],["内容形态偏好","产品图+功能视频+案例","技术白皮书+算法图解","长文观点+行业报告摘要"],["目标受众层级","集成商+工程商渠道","技术决策者+高端渠道","CSO/CISO+企业终端"],["品牌差异化锚点","规模领先+价格优势","研发投入+算法深度","开放生态+隐私合规"],["互动内容策略","评论区行动号召明显","技术问答+专家互动","社区讨论+合规议题引导"],["CTA转化方式","展会预约+产品咨询","技术演示申请+白皮书下载","方案评估+认证合作伙伴"]],"image_theme":"对比分析数据表格"},{"type":"flow_insight","heading":"安防行业内容营销演进路径","banner":{"flow_items":[{"num":"01","desc":"产品参数驱动阶段：硬件规格和价格为核心竞争力，内容以产品图册和功能列表为主"},{"num":"02","desc":"场景解决方案阶段：从产品转向场景叙事，行业案例和部署效果成为主要内容形态"},{"num":"03","desc":"AI技术智能化阶段：算法能力、数据价值和平台生态成为品牌差异化的新战场"},{"num":"04","desc":"合规与价值观阶段：隐私保护、网络安全合规成为高端品牌内容护城河"}]},"cards":[{"title":"头部品牌共同信号","body":"三大品牌均在加速内容专业化升级，减少纯产品推介比例，增加行业洞察和解决方案叙事。AI、IoT、云连接是共同高频词，显示行业叙事已进入智能化阶段的集体共识。"},{"title":"差异化竞争格局","body":"Hikvision以量取胜构建声量护城河，Dahua以技术深度建立专业壁垒，Axis以思想领导力锁定高端决策者，三者已形成清晰的错位竞争态势，直接正面冲突较少。"},{"title":"新兴机会窗口","body":"合规与隐私保护内容严重供给不足，仅Axis有系统布局，中国安防品牌在欧美市场面临合规叙事赤字，这是差异化突破的关键内容空白点，建议优先占位。"}],"image_theme":"行业演进趋势图"},{"type":"metrics","heading":"关键内容指标横向对比","subtitle":"基于领英公开数据与内容分析整理","metrics":[{"value":"周5条","label":"Hikvision频率","trend":"高频多点覆盖","color_idx":0},{"value":"周2条","label":"Dahua频率","trend":"深度精准输出","color_idx":1},{"value":"周1条","label":"Axis频率","trend":"质量优先策略","color_idx":2},{"value":"800字","label":"Axis内容均长","trend":"思想领导力导向","color_idx":0}],"image_theme":"数据仪表盘视觉"},{"type":"section","heading":"核心结论与行动建议","label":"PART 03","image_theme":"策略规划路线图"},{"type":"summary","heading":"核心结论与行动建议","cards":[{"title":"内容策略升级","items":[{"label":"立即行动","body":"将产品推介类内容比例从当前约60%压缩至30%以内，以行业解决方案场景叙事替代，重点围绕客户业务结果而非产品功能参数重构内容框架，参照Axis的思想领导力内容模式。"},{"label":"中期布局","body":"建立垂直行业内容专题库，优先聚焦交通、零售、医疗三个高增长场景，每个场景生产完整的白皮书+案例视频+客户证言三件套内容资产，支撑6-12个月的持续内容输出。"},{"label":"话题策略","body":"主动占领隐私合规与网络安全内容赛道，目前三大竞品中相对薄弱的共同空白，可通过发布行业合规指南和GDPR适用性白皮书，率先建立该领域的内容心智领导权。"}]},{"title":"差异化定位","items":[{"label":"品牌叙事","body":"明确品牌差异化坐标：若与Hikvision竞争需突出技术深度和垂直行业专业度；若与Axis竞争需强化性价比和部署灵活性；避免在所有维度同时竞争，防止品牌形象模糊化。"},{"label":"受众聚焦","body":"从覆盖所有人转向深度触达高价值决策者，建议优先锁定企业安全经理和IT总监决策层，定制化生产解决其工作痛点的实用性内容，系统性降低单位转化成本。"},{"label":"合作生态","body":"学习Axis的生态联合内容策略，主动发起与ERP系统、视频分析软件平台的联合内容合作，通过合作方背书扩大触达范围，同时强化开放生态的品牌定位叙事。"}]},{"title":"执行优先级","items":[{"label":"30天内","body":"完成内容审计：梳理过去90天已发布内容，按照产品推介/解决方案/思想领导力三类重新分类，识别当前内容结构失衡点，制定具体的内容结构调整目标和KPI拆解方案。"},{"label":"90天内","body":"建立行业话题日历：制定年度展会节点内容爆发计划，提前6周启动内容储备，覆盖Intersec、IFSEC三大展会，确保每次展会前后内容密度达到日常3倍以上。"},{"label":"持续追踪","body":"建立竞品内容监控机制，每月产出竞对内容分析报告，持续追踪竞品的话题变化和内容创新，将AI分析能力嵌入常态化竞情监测工作流，保持动态竞争优势。"}]}],"cta":"制定90天内容升级行动计划，以内容力驱动品牌溢价与销售转化的双轮增长","image_theme":"战略规划执行路线"}]}

原文内容如下：
"""


def _strip_lines_noise(text, min_len=15):
    """Extract meaningful content lines from analysis text for fallback/repair."""
    result = []
    for line in text.split('\n'):
        line = line.strip()
        # Remove emoji leaders, markdown markers, separators
        line = re.sub(r'^[\s🔹🔸✅🚀📊⚡✨🎯🔍📡🔥⚖️>#\-\*\|=─━]+', '', line).strip()
        line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)   # bold
        line = re.sub(r'`([^`]+)`', r'\1', line)        # code
        if len(line) >= min_len and not line.startswith('http'):
            result.append(line[:200])  # allow longer lines for content repair pool
    return result


def _fix_val(val, fallback, max_len=200, min_len=4):
    """Return val if it's a real non-empty string, else fallback."""
    BAD = ('待补充','暂无','未知','无','示例','占位','content here',
           'placeholder','this is','lorem','ipsum','...','（暂无）',
           '请补充','内容待填','详细内容','分析内容','具体分析',
           '市场分析','品牌分析','竞对分析','策略建议','核心洞察','数据分析')
    if not val:
        return fallback[:max_len]
    s = str(val).strip()
    if len(s) < min_len or any(b in s.lower() for b in BAD):
        return fallback[:max_len]
    return s[:max_len]


def _fix_body(val, fallback, min_len=40, max_len=200):
    """Stricter validation for body/desc/value content fields.
    Requires real sentence-length content, not just a few words."""
    s = _fix_val(val, fallback, max_len=max_len, min_len=min_len)
    # If still too short after fixing, expand with fallback
    if len(s) < min_len:
        return fallback[:max_len]
    return s


def validate_and_fix_slides(slides, analysis_text, title=''):
    """
    Pass every slide through strict validation.
    Any empty / placeholder field is filled with real extracted content.
    """
    chunks = _strip_lines_noise(analysis_text)
    pool   = list(chunks)   # consumable pool
    ci     = [0]

    def pull(fallback=''):
        """Pull next real chunk from pool, else return fallback."""
        while ci[0] < len(pool):
            val = pool[ci[0]]; ci[0] += 1
            if val and len(val.strip()) >= 15:
                return val.strip()[:150]   # longer chunks = more content in repair
        return fallback or '基于原文深度分析，该品牌在行业中具有明显的差异化定位，通过专业化内容策略持续强化品牌认知，聚焦高价值目标受众，建立清晰的内容到转化的完整路径，形成可持续的品牌竞争壁垒'

    fixed = []
    for s in slides:
        if not isinstance(s, dict) or not s.get('type'):
            continue
        t = s.get('type')

        # ── cover ──────────────────────────────────────────────────────────────
        if t == 'cover':
            s['heading']    = _fix_val(s.get('heading'),    title or '社媒竞对AI分析报告', 40)
            s['subheading'] = _fix_val(s.get('subheading'), '基于多平台内容的深度竞对分析', 60)
            # brands
            brands = s.get('brands', [])
            if not isinstance(brands, list): brands = []
            while len(brands) < 3:
                brands.append({'name': f'品牌{chr(65+len(brands))}', 'role': '分析对象'})
            for b in brands:
                if not isinstance(b, dict): b = {'name':'品牌','role':'分析对象'}
                b['name'] = _fix_val(b.get('name'), '品牌', 30)
                b['role'] = _fix_val(b.get('role'), '分析对象', 12)
            s['brands'] = brands[:3]
            # kpi_badges
            if not s.get('kpi_badges'):
                s['kpi_badges'] = ['AI深度分析', '多维度对比', '数据驱动']

        # ── agenda ─────────────────────────────────────────────────────────────
        elif t == 'agenda':
            s['heading'] = _fix_val(s.get('heading'), '汇报目录', 20)
            items = [i for i in s.get('items', []) if isinstance(i, str) and len(i.strip()) > 2]
            if len(items) < 4:
                items = ['01. 品牌定位与市场卡位', '02. 内容策略深度对比',
                         '03. 传播逻辑与转化路径', '04. 行业洞察与机会', '05. 核心结论与建议']
            s['items'] = items[:7]

        # ── section ────────────────────────────────────────────────────────────
        elif t == 'section':
            s['heading'] = _fix_val(s.get('heading'), pull('章节分析'), 25)
            if not s.get('label') or s.get('label') == 'PART 01':
                # Auto-number sections in order
                section_count = sum(1 for x in fixed if x.get('type') == 'section')
                s['label'] = f'PART {section_count + 1:02d}'

        # ── profile_cards ──────────────────────────────────────────────────────
        elif t == 'profile_cards':
            s['heading'] = _fix_val(s.get('heading'), '品牌核心定位对比', 30)
            entities = s.get('entities', [])
            if not isinstance(entities, list): entities = []
            while len(entities) < 3:
                entities.append({})
            LABELS = ['品牌位势', '核心市场定位', '核心目标市场', '核心受众', '核心商业目标']
            for e in entities[:3]:
                if not isinstance(e, dict): e = {}
                e['name'] = _fix_val(e.get('name'), pull('品牌'), 25)
                e['tag']  = _fix_val(e.get('tag'),  '分析对象', 8)
                fields = e.get('fields', [])
                if not isinstance(fields, list): fields = []
                existing = {f.get('label','') for f in fields if isinstance(f, dict)}
                for lbl in LABELS:
                    if lbl not in existing:
                        fields.append({'label': lbl, 'value': pull(f'该品牌的{lbl}特征分析')})
                for f in fields:
                    if isinstance(f, dict):
                        f['value'] = _fix_body(f.get('value'), pull(f'分析内容'), min_len=30)
                e['fields'] = [f for f in fields if isinstance(f, dict)][:5]
            s['entities'] = entities[:3]

        # ── content_strategy ───────────────────────────────────────────────────
        elif t == 'content_strategy':
            s['heading'] = _fix_val(s.get('heading'), '内容策略对比', 30)
            entities = s.get('entities', [])
            if not isinstance(entities, list): entities = []
            while len(entities) < 3:
                entities.append({})
            for e in entities[:3]:
                if not isinstance(e, dict): e = {}
                e['name']           = _fix_val(e.get('name'), pull('品牌'), 20)
                e['strategy_label'] = _fix_val(e.get('strategy_label'), '核心打法', 14)
                pts = [p for p in e.get('points', []) if isinstance(p, str) and len(p.strip()) > 5]
                while len(pts) < 3:
                    pts.append(pull('基于原文提取的内容策略要点与核心打法分析'))
                e['points'] = [_fix_body(p, pull('内容策略要点'), min_len=40)[:200] for p in pts[:4]]
            s['entities'] = entities[:3]

        # ── dual_panel ─────────────────────────────────────────────────────────
        elif t == 'dual_panel':
            s['heading'] = _fix_val(s.get('heading'), '对比分析', 30)
            panels = s.get('panels', [])
            if not isinstance(panels, list): panels = []
            while len(panels) < 2:
                panels.append({'title': f'分析维度{len(panels)+1}', 'type': 'bullets', 'items': []})
            for p in panels[:2]:
                if not isinstance(p, dict): p = {'title': '维度', 'type': 'bullets', 'items': []}
                p['title'] = _fix_val(p.get('title'), '分析维度', 15)
                items = [it for it in p.get('items', []) if isinstance(it, dict)]
                while len(items) < 3:
                    items.append({'title': f'要点{len(items)+1}', 'desc': pull('基于原文的具体描述内容')})
                for it in items:
                    if not isinstance(it, dict): continue
                    for key in ('title', 'label'):
                        if it.get(key): it[key] = _fix_val(it[key], '要点', 18)
                    it['desc'] = _fix_body(it.get('desc'), pull('具体分析内容'), min_len=25)
                p['items'] = items[:4]
            s['panels'] = panels[:2]

        # ── flow_insight ───────────────────────────────────────────────────────
        elif t == 'flow_insight':
            s['heading'] = _fix_val(s.get('heading'), '市场洞察', 30)
            # banner
            banner = s.get('banner', {})
            if not isinstance(banner, dict): banner = {}
            fi = banner.get('flow_items', [])
            bt = banner.get('text', '')
            if not fi and not bt:
                banner = {'text': pull('行业演进趋势与核心洞察内容')}
            elif fi:
                for item in fi:
                    if isinstance(item, dict):
                        item['desc'] = _fix_val(item.get('desc'), pull('演进阶段'), 30)
            s['banner'] = banner
            # cards
            cards = s.get('cards', [])
            if not isinstance(cards, list): cards = []
            while len(cards) < 3:
                cards.append({'title': f'洞察{len(cards)+1}', 'body': pull('基于原文的市场洞察')})
            for c in cards[:3]:
                if not isinstance(c, dict): c = {}
                c['title'] = _fix_val(c.get('title'), '洞察要点', 15)
                c['body']  = _fix_body(c.get('body'),  pull('基于原文的行业洞察分析内容'), min_len=30)
            s['cards'] = cards[:3]

        # ── three_analysis ─────────────────────────────────────────────────────
        elif t == 'three_analysis':
            s['heading'] = _fix_val(s.get('heading'), '深度分析', 30)
            cards = s.get('cards', [])
            if not isinstance(cards, list): cards = []
            while len(cards) < 3:
                cards.append({'title': f'分析{len(cards)+1}', 'items': []})
            for c in cards[:3]:
                if not isinstance(c, dict): c = {}
                c['title'] = _fix_val(c.get('title'), '分析要点', 15)
                items = c.get('items', [])
                body  = c.get('body', '')
                if not items and not body:
                    # No items and no body — fill body
                    c['body'] = pull('基于原文的分析内容要点')
                elif items:
                    valid_items = [it for it in items if isinstance(it, dict)]
                    while len(valid_items) < 3:
                        valid_items.append({'label': f'要点{len(valid_items)+1}',
                                            'body': pull('分析内容')})
                    for it in valid_items:
                        it['body'] = _fix_body(it.get('body'), pull('具体分析内容'), min_len=25)
                        if it.get('label'):
                            it['label'] = _fix_val(it['label'], '要点', 12)
                    c['items'] = valid_items[:5]
            s['cards'] = cards[:3]

        # ── comparison_table ───────────────────────────────────────────────────
        elif t == 'comparison_table':
            s['heading'] = _fix_val(s.get('heading'), '全景对比总览', 30)
            cols = s.get('cols', ['维度', '品牌A', '品牌B', '品牌C'])
            if not isinstance(cols, list) or len(cols) < 2:
                cols = ['维度', '品牌A', '品牌B', '品牌C']
            n = len(cols)
            rows = s.get('rows', [])
            if not isinstance(rows, list): rows = []
            # Fix each row
            fixed_rows = []
            for row in rows:
                if not isinstance(row, list): continue
                while len(row) < n:
                    row.append(pull('—')[:25])
                fixed_rows.append([_fix_val(cell, pull('分析项'), 30) for cell in row[:n]])
            s['rows'] = fixed_rows
            s['cols'] = cols

        # ── metrics ────────────────────────────────────────────────────────────
        elif t == 'metrics':
            s['heading']  = _fix_val(s.get('heading'), '数据看板', 20)
            s['subtitle'] = _fix_val(s.get('subtitle'), '基于原文数据整理', 35)
            metrics = s.get('metrics', [])
            if not isinstance(metrics, list) or not metrics:
                metrics = [{'value': 'AI分析', 'label': '分析模式', 'trend': '深度', 'color_idx': 0},
                            {'value': '多维', 'label': '对比维度', 'trend': '全面', 'color_idx': 1},
                            {'value': '系统', 'label': '洞察深度', 'trend': '专业', 'color_idx': 2}]
            for m in metrics:
                if not isinstance(m, dict): continue
                m['value'] = _fix_val(m.get('value'), '—', 15)
                m['label'] = _fix_val(m.get('label'), '指标', 10)
            s['metrics'] = metrics[:4]

        # ── bar_chart ──────────────────────────────────────────────────────────
        elif t == 'bar_chart':
            s['heading']  = _fix_val(s.get('heading'), '数据对比', 20)
            s['subtitle'] = _fix_val(s.get('subtitle'), '数据来源：原文整理', 30)
            bars = s.get('bars', [])
            if not isinstance(bars, list): bars = []
            for b in bars:
                if not isinstance(b, dict): continue
                b['label'] = _fix_val(b.get('label'), '对比项', 12)
                if b.get('value') is None: b['value'] = 50
                if b.get('max') is None:   b['max'] = 100
            s['bars'] = bars[:6]

        # ── summary ────────────────────────────────────────────────────────────
        elif t == 'summary':
            s['heading'] = _fix_val(s.get('heading'), '核心结论与行动建议', 25)
            cards = s.get('cards', [])
            points = s.get('points', [])
            if not isinstance(cards, list): cards = []
            if cards:
                while len(cards) < 3:
                    cards.append({'title': f'建议{len(cards)+1}', 'items': []})
                for c in cards[:3]:
                    if not isinstance(c, dict): c = {}
                    c['title'] = _fix_val(c.get('title'), '核心建议', 12)
                    items = c.get('items', [])
                    if not isinstance(items, list): items = []
                    while len(items) < 3:
                        items.append({'label': f'建议{len(items)+1}', 'body': pull('基于原文的可执行建议')})
                    for it in items:
                        if not isinstance(it, dict): continue
                        it['label'] = _fix_val(it.get('label'), '建议', 10)
                        it['body']  = _fix_body(it.get('body'), pull('具体可执行建议'), min_len=30)
                    c['items'] = items[:4]
                s['cards'] = cards[:3]
            elif not points:
                s['points'] = [pull('基于原文的核心结论') for _ in range(5)]

            s['cta'] = _fix_val(s.get('cta'), '基于以上分析，制定针对性内容策略并持续追踪效果', 55)

        fixed.append(s)

    # Guarantee minimum slide set
    types_present = {s.get('type') for s in fixed}
    if 'cover' not in types_present:
        fixed.insert(0, {'type': 'cover', 'heading': title or '竞对分析报告',
                         'subheading': '基于AI深度分析',
                         'brands': [{'name':'品牌A','role':'分析对象'},
                                    {'name':'品牌B','role':'对标品牌'},
                                    {'name':'品牌C','role':'参照'}],
                         'kpi_badges': ['AI深度分析', '多维对比']})
    if 'summary' not in types_present:
        fixed.append({'type': 'summary', 'heading': '核心结论与行动建议',
                      'points': [pull('核心结论') for _ in range(4)],
                      'cta': '结合以上分析，制定针对性内容策略并持续追踪效果。'})

    print(f'[Validate] {len(fixed)} slides validated & fixed', file=sys.stderr)
    return fixed


def _clean_json_string(raw):
    """Multi-layer JSON cleanup for AI-generated strings."""
    raw = raw.strip()
    # Strip markdown fences (```json ... ``` or ``` ... ```)
    raw = re.sub(r'^```[a-zA-Z]*\s*', '', raw, flags=re.MULTILINE)
    raw = re.sub(r'```\s*$', '', raw, flags=re.MULTILINE)
    raw = raw.strip()
    # Extract from first { to last }
    start = raw.find('{')
    end   = raw.rfind('}') + 1
    if start < 0 or end <= start:
        return None
    return raw[start:end]


def _parse_json_robust(json_str):
    """Try parsing JSON, then attempt auto-repair for common AI mistakes."""
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass
    # Fix 1: trailing commas before } or ]
    s = re.sub(r',\s*([}\]])', r'\1', json_str)
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        pass
    # Fix 2: single-quoted strings → double-quoted (careful not to break apostrophes)
    s2 = re.sub(r"(?<=[{,:\[])\s*'([^']*?)'", r'"\1"', s)
    try:
        return json.loads(s2)
    except json.JSONDecodeError:
        pass
    # Fix 3: unquoted keys → quoted keys
    s3 = re.sub(r'([{,]\s*)([a-zA-Z_][a-zA-Z0-9_]*)(\s*:)', r'\1"\2"\3', s2)
    try:
        return json.loads(s3)
    except Exception:
        pass
    return None


def structure_analysis(api_key, text_model, analysis_text, title, subtitle, proxy_url=''):
    prompt = STRUCTURE_PROMPT + '\n\n' + analysis_text
    messages = [
        {'role': 'system',
         'content': ('你是专业商务PPT结构化AI。只输出纯JSON，第一个字符必须是{，'
                     '最后一个字符必须是}，禁止任何非JSON内容、markdown标记、解释性文字。'
                     '所有字段必须有实际内容，禁止空字符串。')},
        {'role': 'user', 'content': prompt},
    ]
    raw = doubao_chat(api_key, text_model, messages, proxy_url, max_tokens=12000)  # generous for 14-18 slides

    json_str = _clean_json_string(raw)
    if not json_str:
        raise ValueError(f'AI未返回有效JSON结构。原始输出前200字: {raw[:200]}')

    data = _parse_json_robust(json_str)
    if data is None:
        raise ValueError(f'JSON解析失败（多种修复均无效）。原始JSON前300字:\n{json_str[:300]}')

    slides = data.get('slides', [])
    if not slides:
        raise ValueError('AI返回的JSON中slides数组为空，请检查提示词或重试')

    # Inject job-level title override
    for s in slides:
        if isinstance(s, dict) and s.get('type') == 'cover':
            if title:    s['heading']    = title
            if subtitle: s['subheading'] = subtitle
            break

    # Full validation + auto-fix pass
    slides = validate_and_fix_slides(slides, analysis_text, title)
    print(f'[Structure] ✅ Done: {len(slides)} slides', file=sys.stderr)
    return slides


def fallback_slides(analysis_text, title, subtitle):
    """Robust fallback: extract real sentences from analysis_text to fill slides."""
    chunks = _strip_lines_noise(analysis_text, min_len=20)
    pool = list(chunks)
    ci = [0]

    def pull(fb='分析内容'):
        while ci[0] < len(pool):
            v = pool[ci[0]]; ci[0] += 1
            if v and len(v) >= 10: return v[:80]
        return fb

    # Extract markdown headings for agenda
    headings = re.findall(r'^#{1,3}\s+(.{3,40})', analysis_text, re.MULTILINE)
    clean_h  = [re.sub(r'[🔹🔸✅🚀📊⚡✨🎯🔍📡🔥⚖️>#*\-|]', '', h).strip()[:30]
                for h in headings if h.strip()][:6]

    def pull3():
        return [pull(f'分析要点{i+1}') for i in range(3)]

    slides = [
        {'type': 'cover',
         'heading':    title or '社媒AI竞对分析报告',
         'subheading': subtitle or '基于多平台内容的深度竞对分析',
         'brands': [{'name':'品牌A','role':'分析对象'},
                    {'name':'品牌B','role':'对标品牌'},
                    {'name':'品牌C','role':'竞争参照'}],
         'kpi_badges': ['AI深度分析', '多维对比', '数据驱动']},

        {'type': 'agenda', 'heading': '汇报目录',
         'items': [f'{i+1:02d}. {h}' for i, h in enumerate(clean_h)] if len(clean_h) >= 3
                  else ['01. 品牌定位与市场卡位', '02. 内容策略深度对比',
                        '03. 传播逻辑分析', '04. 行业洞察与机会', '05. 核心结论与建议']},

        {'type': 'section', 'heading': '品牌定位与内容策略', 'label': 'PART 01'},

        {'type': 'three_analysis', 'heading': '核心分析洞察',
         'cards': [
             {'title': '品牌定位', 'items': [{'label': f'要点{i+1}', 'body': pull('品牌分析')} for i in range(4)]},
             {'title': '内容策略', 'items': [{'label': f'要点{i+1}', 'body': pull('内容分析')} for i in range(4)]},
             {'title': '传播逻辑', 'items': [{'label': f'要点{i+1}', 'body': pull('传播分析')} for i in range(4)]},
         ]},

        {'type': 'section', 'heading': '洞察与建议', 'label': 'PART 02'},

        {'type': 'three_analysis', 'heading': '市场洞察与机会',
         'cards': [
             {'title': '行业趋势', 'body': pull('行业趋势分析')},
             {'title': '竞争格局', 'body': pull('竞争格局分析')},
             {'title': '市场机会', 'body': pull('市场机会分析')},
         ]},

        {'type': 'summary', 'heading': '核心结论与行动建议',
         'points': [pull('结论') for _ in range(5)],
         'cta':    '结合以上分析，制定针对性内容策略并持续追踪效果。'},
    ]
    slides = validate_and_fix_slides(slides, analysis_text, title)
    print(f'[Fallback] {len(slides)} slides built from real content', file=sys.stderr)
    return slides



# ═══════════════════════════════════════════════════════════════════════════════
# Step 2a – LLM generates rich image prompts for each visual slide
# ═══════════════════════════════════════════════════════════════════════════════

IMAGE_PROMPT_SYSTEM = """You are a professional AI image prompt engineer specializing in Doubao Seedream visual model prompts for business PPT slides.

STRICT RULES — violating any rule means FAILED output:
1. English only, 150-220 words, rich and specific visual description, no Chinese, no explanation, no label, output the prompt directly — never cut short mid-sentence
2. HARD REQUIREMENT: absolutely no text, no letters, no numbers, no watermarks, no logos, no signs, no banners — any text in the image = FAILED
3. HARD REQUIREMENT: no human faces, no identifiable people — abstract or blurred figures only if needed
4. Style: premium professional photography OR cinematic 3D render, dark navy/charcoal background (#111827 tone), sophisticated commercial quality
5. Lighting: soft blue-tinted professional studio lighting, high contrast highlights, clean sharp focus, 8K ultra HD
6. Composition rules by slide type:
   - Cover / Section: full-frame centered composition, dramatic depth of field, hero visual
   - Content pages: LEFT or RIGHT third of frame intentionally minimal/dark (reserved for PPT text overlay), main subject fills other 2/3
   - All pages: 16:9 widescreen, no element touching frame edges
7. Color palette: dark navy blue (#111827), electric blue accents (#3B82F6), emerald green accents (#10B981), no bright saturated colors, no warm oranges/reds
8. Forbidden elements: any text, letters, numbers, watermarks, logos, icons, faces, hands pointing, explicit business clichés (handshakes, etc.)

Output: the English prompt only, starting directly with the visual description."""

def generate_image_prompts(api_key, text_model, slides, proxy_url='', status_path=None):
    """Step 2a: LLM generates tailored image prompts for visual slides."""
    # All slide types that benefit from a background image
    visual_types = {'cover', 'section', 'content_strategy', 'flow_insight', 'dual_panel', 'three_analysis', 'summary', 'profile_cards', 'comparison_table', 'metrics'}
    needs_image  = [i for i,s in enumerate(slides) if s.get('type') in visual_types]

    if not needs_image:
        return slides

    print(f'[ImgPrompts] Generating prompts for {len(needs_image)} slides...', file=sys.stderr)

    for idx in needs_image:
        s    = slides[idx]
        stype = s.get('type')

        # Build a short context summary for the LLM
        ctx_parts = [f"幻灯片类型: {stype}"]
        if s.get('heading'):   ctx_parts.append(f"标题: {s['heading']}")
        if s.get('subheading'):ctx_parts.append(f"副标题: {s['subheading']}")

        if stype == 'content_strategy':
            entities = s.get('entities',[])
            names = [e.get('name','') for e in entities[:3] if isinstance(e,dict)]
            strategies = [e.get('strategy_label','') for e in entities[:3] if isinstance(e,dict)]
            ctx_parts.append(f"品牌: {', '.join(names)}")
            ctx_parts.append(f"内容策略: {'; '.join(strategies)}")
        elif stype == 'flow_insight':
            cards = s.get('cards',[])
            titles = [c.get('title','') for c in cards[:3] if isinstance(c,dict)]
            ctx_parts.append(f"洞察: {'; '.join(titles)}")
        elif stype == 'dual_panel':
            panels = s.get('panels',[])
            panel_titles = [p.get('title','') for p in panels if isinstance(p,dict)]
            ctx_parts.append(f"内容: {'; '.join(panel_titles)}")

        if s.get('image_theme'):
            ctx_parts.append(f"图像主题关键词: {s['image_theme']}")
        context = '\n'.join(ctx_parts)
        # Determine composition hint based on slide type
        comp_hint = {
            'cover':            'full-frame centered hero composition — dramatic and impressive',
            'section':          'full-frame centered composition with strong visual depth',
            'content_strategy': 'right third dark/minimal for text overlay, main subject fills left 2/3',
            'flow_insight':     'lower third dark/minimal for card overlay, upper 2/3 atmospheric',
            'dual_panel':       'center subject, both sides have minimal dark space for text panels',
            'profile_cards':    'right third intentionally minimal for text overlay, left 2/3 abstract corporate architecture',
            'comparison_table': 'full-frame dark abstract background, subtle grid or matrix pattern, no bright focal point',
            'metrics':          'center composition, abstract data flow lines or glowing circuits, dark background with accent lights',
        }.get(stype, 'right third intentionally minimal for text overlay, subject in left 2/3')

        user_msg = f"""Generate a high-quality English image prompt for Doubao Seedream for this PPT slide.

SLIDE INFO:
{context}

COMPOSITION REQUIREMENT: {comp_hint}

CRITICAL: output must start directly with the visual description. No Chinese. No text in image. No faces."""

        try:
            messages = [
                {'role':'system','content':IMAGE_PROMPT_SYSTEM},
                {'role':'user','content':user_msg},
            ]
            prompt = doubao_chat(api_key, text_model, messages, proxy_url, max_tokens=600)
            prompt = prompt.strip().strip('"').strip("'")
            slides[idx]['image_prompt'] = prompt
            print(f'  [Prompt {idx+1}] ({len(prompt)} chars) {prompt}', file=sys.stderr)
        except Exception as e:
            print(f'  [Prompt {idx+1}] Error: {e}', file=sys.stderr)
            # Fallback prompt — strict no-text, dark theme
            topic = s.get('heading','corporate strategy')[:40]
            slides[idx]['image_prompt'] = (
                f"Premium business technology photography, {topic} concept, abstract geometric patterns, "
                "dark navy blue background (#111827), electric blue accent lighting, no text whatsoever, "
                "absolutely no letters or numbers, no human faces, minimalist sophisticated composition, "
                "right third of frame intentionally dark for text overlay, 8K ultra HD, 16:9 widescreen, "
                "commercial photography style, high contrast clean sharp focus."
            )

    return slides


# ═══════════════════════════════════════════════════════════════════════════════
# Step 2b – Visual model generates images from LLM prompts
# ═══════════════════════════════════════════════════════════════════════════════

def generate_slide_images(api_key, slides, img_dir, proxy_url='', status_path=None, base_progress=30):
    """Step 2b: Seedream visual model generates images from Step-2a LLM prompts."""
    os.makedirs(img_dir, exist_ok=True)
    img_map = {}
    # Include all slide types that benefit from imagery
    visual_types = {'content_strategy', 'cover', 'section', 'flow_insight', 'dual_panel'}
    visual_slides = [(i, s) for i, s in enumerate(slides)
                     if s.get('type') in visual_types and s.get('image_prompt', '').strip()]
    total_v = len(visual_slides)
    if total_v == 0:
        return img_map

    for vi, (i, slide) in enumerate(visual_slides):
        prompt = slide.get('image_prompt', '').strip()
        progress = int(base_progress + (vi / max(total_v, 1)) * 38)
        if status_path:
            write_status(status_path, 'running', 2,
                         f'🎨 视觉大模型生图 {vi+1}/{total_v}（{IMG_MODEL}）...', progress)

        print(f'[Seedream] Calling for slide type={slide.get("type")} prompt_len={len(prompt)}', file=sys.stderr)
        img_bytes = seedream_generate(api_key, prompt, proxy_url)
        if img_bytes:
            if HAS_PIL:
                try:
                    img = PILImage.open(io.BytesIO(img_bytes)).convert('RGB')
                    img = img.resize((1280, 720), PILImage.LANCZOS)
                    buf = io.BytesIO()
                    img.save(buf, format='JPEG', quality=90)
                    img_bytes = buf.getvalue()
                except Exception:
                    pass
            out_path = os.path.join(img_dir, f'slide_{i}.jpg')
            with open(out_path, 'wb') as f:
                f.write(img_bytes)
            img_map[i] = out_path
            print(f'  [Seedream] slide {i+1} ✅ saved', file=sys.stderr)
        else:
            print(f'  [Seedream] slide {i+1} ❌ failed', file=sys.stderr)
        time.sleep(1.0)

    return img_map

def assemble_ppt(slides_data, img_map, output_path, date_str=''):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    total = len(slides_data) + 1

    for i, s in enumerate(slides_data):
        stype = s.get('type','three_analysis')
        num   = i + 1

        if stype == 'cover':
            s2 = dict(s); s2['_img'] = img_map.get(i)
            make_cover(prs, s2, date_str)
        elif stype == 'agenda':
            make_agenda(prs, s, num, total)
        elif stype == 'section':
            s2 = dict(s); s2['_img'] = img_map.get(i)
            make_section(prs, s2, num, total)
        elif stype == 'profile_cards':
            make_profile_cards(prs, s, num, total)
        elif stype == 'content_strategy':
            s2 = dict(s); s2['_img_map'] = {j: img_map.get(i+j) for j in range(3) if img_map.get(i+j)}
            make_content_strategy(prs, s2, num, total)
        elif stype == 'dual_panel':
            s2 = dict(s); s2['_img'] = img_map.get(i)
            make_dual_panel(prs, s2, num, total)
        elif stype == 'flow_insight':
            s2 = dict(s); s2['_img'] = img_map.get(i)
            make_flow_insight(prs, s2, num, total)
        elif stype == 'three_analysis':
            s2 = dict(s); s2['_img'] = img_map.get(i)
            make_three_analysis(prs, s2, num, total)
        elif stype == 'comparison_table':
            make_comparison_table(prs, s, num, total)
        elif stype == 'metrics':
            make_metrics(prs, s, num, total)
        elif stype == 'bar_chart':
            make_bar_chart(prs, s, num, total)
        elif stype == 'summary':
            make_summary(prs, s, num, total)
        else:
            # Legacy fallback
            make_three_analysis(prs, s, num, total)

    make_end(prs)
    prs.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# Main pipeline
# ═══════════════════════════════════════════════════════════════════════════════

def run(job):
    api_key       = job['api_key']
    proxy_url     = job.get('proxy_url','')
    analysis_text = job['analysis_text']
    title         = job.get('title','社媒AI竞对分析报告')
    subtitle      = job.get('subtitle','')
    date_str      = job.get('date','')
    output_path   = job['output_path']
    status_path   = job['status_path']
    enable_images = job.get('enable_images', True)
    text_model    = job.get('text_model','doubao-seed-2-0-pro-260215')

    # Allow job to override the visual model (user-configurable in settings)
    global IMG_MODEL
    if job.get('visual_model'):
        IMG_MODEL = job['visual_model']

    print(f'[PPT] TEXT_MODEL={text_model}  IMG_MODEL={IMG_MODEL}', file=sys.stderr)

    img_dir = output_path.replace('.pptx','_imgs')
    os.makedirs(img_dir, exist_ok=True)

    write_status(status_path,'running',1,f'🤖 Step 1/3: {text_model} 结构化分析中...',5)
    try:
        slides = structure_analysis(api_key, text_model, analysis_text, title, subtitle, proxy_url)
        write_status(status_path,'running',1,f'✅ 结构化完成，共{len(slides)}张幻灯片',28)
    except Exception as e:
        print(f'[Structure] fallback: {e}', file=sys.stderr)
        slides = fallback_slides(analysis_text, title, subtitle)
        write_status(status_path,'running',1,f'⚠️ 已用备用结构（{len(slides)}张）',28)

    img_map = {}
    if enable_images and api_key:
        # Step 2a: LLM generates tailored image prompts
        write_status(status_path,'running',2,f'🤖 Step 2a/3: {text_model} 生成配图提示词...',30)
        try:
            slides = generate_image_prompts(api_key, text_model, slides, proxy_url, status_path)
            n_prompted = sum(1 for s in slides if s.get('image_prompt','').strip())
            write_status(status_path,'running',2,f'✅ 提示词生成完成（{n_prompted}张幻灯片）',38)
        except Exception as e:
            print(f'[ImgPrompts] {e}', file=sys.stderr)
            write_status(status_path,'running',2,'⚠️ 提示词生成问题，使用默认提示词',38)

        # Step 2b: Seedream visual model generates images
        write_status(status_path,'running',2,f'🎨 Step 2b/3: {IMG_MODEL} 生成配图...',40)
        try:
            img_map = generate_slide_images(api_key, slides, img_dir, proxy_url, status_path, 40)
            write_status(status_path,'running',2,f'✅ 配图完成（{len(img_map)}张）',70)
        except Exception as e:
            print(f'[Images] {e}', file=sys.stderr)
            write_status(status_path,'running',2,'⚠️ 配图生成问题，继续无图版',70)
    else:
        write_status(status_path,'running',2,'⏭️ 跳过配图（纯矢量模式）',70)

    write_status(status_path,'running',3,'📄 Step 3/3: 组装 PPT 文件...',73)
    assemble_ppt(slides, img_map, output_path, date_str)
    size_kb = os.path.getsize(output_path) // 1024
    write_status(status_path,'running',3,
                 f'✅ PPT已生成（{size_kb}KB，{len(slides)+1}页，{text_model}）',95)

    # Save preview JSON
    job_id = job.get('job_id','')
    if job_id:
        preview_path = os.path.join(os.path.dirname(status_path), f'{job_id}-slides.json')
        try:
            preview_slides = []
            for i, s in enumerate(slides):
                ps = {k:v for k,v in s.items() if k not in ('image_prompt','_img_map')}
                img_path = img_map.get(i)
                if img_path and os.path.exists(img_path):
                    with open(img_path,'rb') as f:
                        ps['img_b64'] = base64.b64encode(f.read()).decode('ascii')
                preview_slides.append(ps)
            with open(preview_path,'w',encoding='utf-8') as f:
                json.dump({'slides':preview_slides,'slide_count':len(slides)+1,
                           'model':text_model,'date':date_str},f,ensure_ascii=False)
        except Exception as e:
            print(f'[Preview] {e}', file=sys.stderr)

    try:
        import shutil; shutil.rmtree(img_dir, ignore_errors=True)
    except: pass

    return output_path, len(slides) + 1


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python3 ppt_pipeline.py <job_json_path>'); sys.exit(1)
    with open(sys.argv[1],'r',encoding='utf-8') as f:
        job = json.load(f)
    status_path = job['status_path']
    try:
        out, slide_count = run(job)
        filename = os.path.basename(out)
        write_status(status_path,'done',4,'🎉 PPT生成完成！',100,url=f'/ppt/{filename}')
        print(f'OK:{out}')
    except Exception as e:
        import traceback; tb = traceback.format_exc()
        print(f'[ERROR] {e}\n{tb}', file=sys.stderr)
        write_status(status_path,'error',0,f'生成失败: {str(e)}',0,error=str(e))
        sys.exit(1)
